"""
Gradient-based optimization of text embeddings for image generation using SDXL.
Uses Adam optimizer to modify text embeddings while maximizing aesthetic and CLIP scores.
"""

# System imports
import sys
import os
import shutil

# Add parent directory to Python path for module imports
parent_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
sys.path.insert(0, parent_dir)

# External imports - consider grouping by functionality
import yaml 
import torch
import torch.nn.functional as F
import numpy as np
import pandas as pd
from diffusers import StableDiffusionXLPipeline, AutoencoderKL
import random
from PIL import Image 
import matplotlib.pyplot as plt
import time
from datasets import load_dataset
import clip
from collections import defaultdict
import csv
from pptx import Presentation
from pptx.util import Inches
import argparse

# Add argument parsing
parser = argparse.ArgumentParser(description='Run optimization with configuration file')
parser.add_argument('--config', type=str, default="algorithms/config/config_adam.yaml",
                   help='Path to configuration YAML file')
args = parser.parse_args()

# Use the provided config path or default
config_path = args.config

with open(config_path, 'r') as file:
    config = yaml.safe_load(file)

# Extract parameters from the configuration
SEED_PATH = config['seed_path']
SEED = config['seed']
cuda_n = str(config['cuda'])
predictor = config['predictor']
num_inference_steps = config['num_inference_steps']
height = config['height']
width = config['width']
NUM_ITERATIONS = config['num_iterations']
model_id = config['model_id']
OUTPUT_FOLDER = config['results_folder']
alpha = config['alpha']
beta = config['beta']
max_aesthetic = config['max_aesthetic']
max_clip = config['max_clip']
adam_lr = float(config['adam_lr'])
adam_weight_decay = float(config['adam_weight_decay'])
adam_eps = float(config['adam_eps'])
adam_beta1 = float(config['adam_beta1'])
adam_beta2 = float(config['adam_beta2'])

if predictor == 0:
    predictor_name = 'simulacra'
elif predictor == 1:
    predictor_name = 'laionv1'
elif predictor == 2:
    predictor_name = 'laionv2'
else:
    raise ValueError("Invalid predictor option.")

OUTPUT_FOLDER = f"{OUTPUT_FOLDER}/adam_clip_{predictor_name}_sdxlturbo_{SEED}_a{int(alpha*100)}_b{int(beta*100)}"

# Save the loaded config to a txt file in the output folder
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
# Copy the YAML configuration file to the output folder
shutil.copy(config_path, os.path.join(OUTPUT_FOLDER, "config_used.yaml"))

# Check if a GPU is available and if not, use the CPU
device = "cuda:" + cuda_n if torch.cuda.is_available() else "cpu"

# Load the Stable Diffusion pipeline
pipe = StableDiffusionXLPipeline.from_pretrained(
    model_id,
    torch_dtype=torch.float32,
    use_safetensors=True,
).to(device)

# Enable gradient computation for the pipeline
pipe.set_progress_bar_config(disable=True)
call_with_grad = pipe.__class__.__call__.__wrapped__.__get__(pipe, pipe.__class__)

# Add warning about computational overhead
#torch.autograd.set_detect_anomaly(True)  # Warning: This may impact performance

# CLIP setup
clip_model_name = "ViT-L/14"
clip_model, clip_preprocess = clip.load(clip_model_name, device=device)

# Prompt dataset loading and preprocessing
prompt_dataset = load_dataset("nateraw/parti-prompts")["train"]

N_PER_CATEGORY = config['prompt_per_categorie']  # Number of prompts to sample per category
SUBSET_SEED = config['prompt_sample_seed']
random.seed(SUBSET_SEED)

# Group prompts by category
category_prompts = defaultdict(list)
for item in prompt_dataset:
    category = item.get("Category", "Uncategorized")
    category_prompts[category].append(item["Prompt"])

# Sample N_PER_CATEGORY prompts from each category and keep track of category
selected_prompts_with_category = []
for category, prompts in category_prompts.items():
    if len(prompts) >= N_PER_CATEGORY:
        sampled = random.sample(prompts, N_PER_CATEGORY)
    else:
        sampled = prompts  # If not enough, take all
    for prompt in sampled:
        selected_prompts_with_category.append((prompt, category))

print(f"Selected {len(selected_prompts_with_category)} prompts from {len(category_prompts)} categories.")

prompt_list_path = os.path.join(OUTPUT_FOLDER, "selected_prompts.txt")
with open(prompt_list_path, "w", encoding="utf-8") as f:
    for prompt, category in selected_prompts_with_category:
        f.write(f"{category}\t{prompt}\n")
print(f"Saved selected prompts to {prompt_list_path}")

# Initialize the aesthetic model
if predictor == 0:
    from aesthetic_evaluation.src import simulacra_rank_image
    aesthetic_model = simulacra_rank_image.SimulacraAesthetic(device)
    model_name = 'SAM'
elif predictor == 1:
    from aesthetic_evaluation.src import laion_rank_image
    aesthetic_model = laion_rank_image.LAIONAesthetic(device, clip_model=clip_model_name)
    model_name = 'LAIONV1'
elif predictor == 2:
    from aesthetic_evaluation.src import laion_v2_rank_image
    aesthetic_model = laion_v2_rank_image.LAIONV2Aesthetic(device, clip_model=clip_model_name)
    model_name = 'LAIONV2'
else:
    raise ValueError("Invalid predictor option.")

if SEED_PATH is None:
    seed_list = [SEED]
else:
    with open(SEED_PATH, 'r') as file:
        seed_list = [int(line.strip()) for line in file]

def generate_image_from_embeddings(text_embeddings, seed):
    generator = torch.Generator(device=device).manual_seed(seed)

    prompt_embeds = text_embeddings[0]
    pooled_prompt_embeds = text_embeddings[1]

    out = call_with_grad(
        prompt_embeds              = prompt_embeds,
        pooled_prompt_embeds       = pooled_prompt_embeds,
        guidance_scale             = 0.0,
        num_inference_steps        = num_inference_steps,
        generator                  = generator,
        height                     = height,
        width                      = width,
        output_type               = "pt"
    )["images"]

    image = out.clamp(0, 1).squeeze(0).permute(1, 2, 0)      # HWC
    return image.to(device)

def aesthetic_evaluation(image):
    # image is a tensor of shape [H, W, C]
    image_input = image.permute(2, 0, 1).to(torch.float32)  # [1, C, H, W]

    if predictor == 0:
        # Simulacra Aesthetic Model
        score = aesthetic_model.predict_from_tensor(image_input)
    elif predictor == 1 or predictor == 2:
        # LAION Aesthetic Predictor V1 and V2
        score = aesthetic_model.predict_from_tensor(image_input)
    else:
        return torch.tensor(0.0, device=device)

    return score

# CLIP (ViT-L/14) image norm
_CLIP_MEAN = torch.tensor([0.48145466, 0.4578275, 0.40821073]).view(1,3,1,1)
_CLIP_STD  = torch.tensor([0.26862954, 0.26130258, 0.27577711]).view(1,3,1,1)
def evaluate_clip_score(image_tensor, text_features):
    """
    image_tensor: [H, W, C], float in [0,1], requires_grad=True
    returns: scalar similarity (higher = better)
    """

    clip_model.eval()

    # --- differentiable preprocess (no PIL, no .byte) ---
    img = image_tensor.permute(2,0,1).unsqueeze(0)            # [1,C,H,W]
    img = img.to(device=device, dtype=torch.float32)
    img = F.interpolate(img, size=(224,224), mode="bicubic", align_corners=False)
    mean = _CLIP_MEAN.to(img.device, img.dtype)
    std  = _CLIP_STD.to(img.device, img.dtype)
    img = (img - mean) / std

    # Encode image WITH grad (through CLIP image tower)
    image_features = clip_model.encode_image(img).float()
    image_features = F.normalize(image_features, dim=-1, eps=1e-6)

    sim = (image_features @ text_features.T).squeeze()  # scalar
    return sim

def format_time(seconds):
    seconds = int(seconds)
    hours = seconds // 3600
    minutes = (seconds % 3600) // 60
    seconds = seconds % 60
    if hours > 0:
        return f"{hours}h {minutes}m {seconds}s"
    elif minutes > 0:
        return f"{minutes}m {seconds}s"
    else:
        return f"{seconds}s"

def main(seed, seed_number, selected_prompt, category, prompt_number):
    torch.manual_seed(seed)
    np.random.seed(seed)
    random.seed(seed)

    print(f"Selected prompt: {selected_prompt} (Category: {category})")

    results_folder = f"{OUTPUT_FOLDER}/results_{model_name}_{seed}_{prompt_number}"
    os.makedirs(results_folder, exist_ok=True)

    # Text features don't depend on your params; compute w/o grad
    with torch.no_grad():
        text_tokens = clip.tokenize([selected_prompt]).to(device)
        text_features = clip_model.encode_text(text_tokens).float()
        text_features = F.normalize(text_features, dim=-1, eps=1e-6)

    (prompt_embeds,
        negative_prompt_embeds,
        pooled_prompt_embeds,
        negative_pooled_prompt_embeds) = pipe.encode_prompt(
                selected_prompt,
                negative_prompt="",                 
                device=device,
                num_images_per_prompt=1,
                do_classifier_free_guidance=False
    )
    text_embeddings_init = [prompt_embeds.detach().clone(), pooled_prompt_embeds.detach().clone()]
    text_embeddings = [torch.nn.Parameter(prompt_embeds.clone()), torch.nn.Parameter(pooled_prompt_embeds.clone())]

    with torch.no_grad():
        initial_image = generate_image_from_embeddings(text_embeddings_init, seed)
        image_np = initial_image.detach().clone().cpu().numpy()
        image_np = (image_np * 255).astype(np.uint8)
        pil_image = Image.fromarray(image_np)
        pil_image.save(f"{results_folder}/it_0.png")

    aesthetic_score = aesthetic_evaluation(initial_image)

    clip_score = evaluate_clip_score(initial_image, text_features)

    initial_combined_score = alpha * aesthetic_score / max_aesthetic + beta * clip_score / max_clip
    initial_combined_loss = 1 - initial_combined_score

    combined_score_list = [initial_combined_score.item()]
    combined_loss_list = [initial_combined_loss.item()]
    time_list = [0]
    best_score = initial_combined_score
    best_text_embeddings = text_embeddings_init.copy()

    optimizer = torch.optim.Adam(text_embeddings, lr=adam_lr, betas=(adam_beta1,adam_beta2), weight_decay=adam_weight_decay, eps=adam_eps)  

    start_time = time.time()

    # Add lists to store the metrics
    aesthetic_score_list = [aesthetic_score.item()]
    clip_score_list = [clip_score.item()]

    for iteration in range(1, NUM_ITERATIONS + 1):
        print(f"Iteration {iteration}/{NUM_ITERATIONS}")

        optimizer.zero_grad()

        #with torch.autocast(device_type=device, dtype=torch.float16):
        image = generate_image_from_embeddings(text_embeddings, seed)
        aesthetic_score = aesthetic_evaluation(image)
        clip_score = evaluate_clip_score(image, text_features)
        combined_score = alpha * aesthetic_score / max_aesthetic + beta * clip_score / max_clip
        combined_loss = 1 - combined_score

        # Calculate gradients
        combined_loss.backward()
        # Update parameters
        optimizer.step()

        # Append metrics to their respective lists
        aesthetic_score_list.append(aesthetic_score.item())
        clip_score_list.append(clip_score.item())

        if combined_score.item() > best_score:
            best_score = combined_score.item()
            best_text_embeddings = text_embeddings.copy()

        combined_score_list.append(combined_score.item())
        combined_loss_list.append(combined_loss.item())

        image_np = image.detach().clone().cpu().numpy()
        image_np = (image_np * 255).astype(np.uint8)
        pil_image = Image.fromarray(image_np)
        pil_image.save(f"{results_folder}/it_{iteration}.png")

        elapsed_time = time.time() - start_time
        iterations_done = iteration
        iterations_left = NUM_ITERATIONS - iteration
        average_time_per_iteration = elapsed_time / iterations_done
        estimated_time_remaining = average_time_per_iteration * iterations_left

        formatted_time_remaining = format_time(estimated_time_remaining)

        time_list.append(elapsed_time)

        # Save metrics to the results DataFrame
        results = pd.DataFrame({
            "iteration": list(range(0, iteration + 1)),
            "prompt": [selected_prompt] + [''] * iteration,
            "category": [category] + [''] * iteration,
            "combined_score": combined_score_list,
            "combined_loss": combined_loss_list,
            "aesthetic_score": aesthetic_score_list,
            "clip_score": clip_score_list,
            "elapsed_time": time_list
        })

        results.to_csv(f"{results_folder}/score_results.csv", index=False, na_rep='nan')

        # Plot and save the fitness evolution
        plot_results(results, results_folder)

        # Print stats
        print(f"Seed {seed_number} Iteration {iteration}/{NUM_ITERATIONS}: Combined Score: {combined_score.item()}, Aesthetic Score: {aesthetic_score.item()}, CLIP Score: {clip_score.item()}, Estimated time remaining: {formatted_time_remaining}")

    # Save the overall best image
    with torch.no_grad():
        best_image = generate_image_from_embeddings(best_text_embeddings, seed)
    best_image_np = best_image.detach().cpu().numpy()
    best_image_np = (best_image_np * 255).astype(np.uint8)
    pil_image = Image.fromarray(best_image_np)
    pil_image.save(f"{results_folder}/best_all.png")

def plot_results(results, results_folder):
    plt.figure(figsize=(10, 6))  # Increase figure size
    plt.plot(results['iteration'], results['aesthetic_score'], label="Aesthetic Score")
    plt.xlabel('Iteration')
    plt.ylabel('Aesthetic Score')
    plt.title('Aesthetic Score Evolution')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))  # Move legend outside the plot
    plt.tight_layout()  # Adjust layout
    plt.savefig(results_folder + "/aesthetic_evolution.png")
    plt.close()

    plt.figure(figsize=(10, 6))  # Increase figure size
    plt.plot(results['iteration'], results['clip_score'], label="CLIP Score")
    plt.xlabel('Iteration')
    plt.ylabel('CLIP Score')
    plt.title('CLIP Score Evolution')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))  # Move legend outside the plot
    plt.tight_layout()  # Adjust layout
    plt.savefig(results_folder + "/clip_evolution.png")
    plt.close()

    # Plot all losses in one plot
    plt.figure(figsize=(10, 6))  # Increase figure size
    plt.plot(results['iteration'], results['combined_loss'], label="Combined Loss")
    plt.xlabel('Iteration')
    plt.ylabel('Loss')
    plt.title('Loss Evolution')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))  # Move legend outside the plot
    plt.tight_layout()  # Adjust layout
    plt.savefig(results_folder + "/loss_evolution.png")
    plt.close()

def plot_mean_std(x_axis, m_vec, std_vec, description, title=None, y_label=None, x_label=None):
    lower_bound = [M_new - Sigma for M_new, Sigma in zip(m_vec, std_vec)]
    upper_bound = [M_new + Sigma for M_new, Sigma in zip(m_vec, std_vec)]

    plt.plot(x_axis, m_vec, '--', label=description + " Avg.")
    plt.fill_between(x_axis, lower_bound, upper_bound, alpha=.3, label=description + " Avg. ± SD")
    if title is not None:
        plt.title(title)
    if y_label is not None:
        plt.ylabel(y_label)
    if x_label is not None:
        plt.xlabel(x_label)

def aggregate_results():
    # Initialize aggregated_data as None
    aggregated_data = None

    # Variables to track the maximum final combined score
    max_final_combined_score = float('-inf')
    max_seed = None
    max_prompt_number = None

    # Iterate over all subdirectories
    for folder_name in os.listdir(OUTPUT_FOLDER):
        if folder_name.startswith(f"results_{model_name}_"):
            seed = folder_name.split("_")[-2]  # Extract the seed number
            prompt_number = folder_name.split("_")[-1]  # Extract the prompt number

            file_path = os.path.join(OUTPUT_FOLDER, folder_name, "score_results.csv")

            if os.path.exists(file_path):
                # Read the CSV file
                df = pd.read_csv(file_path)

                df = df.drop(columns=["prompt", "category"]) 

                df = df.rename(columns={"combined_score": f"combined_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"combined_loss": f"combined_loss_{seed}_{prompt_number}"})
                df = df.rename(columns={"aesthetic_score": f"aesthetic_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"clip_score": f"clip_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"elapsed_time": f"elapsed_time_{seed}_{prompt_number}"})

                # Check the final combined score in this file
                max_score = df[f"combined_score_{seed}_{prompt_number}"].max()
                if max_score > max_final_combined_score:
                    max_final_combined_score = max_score
                    max_seed = seed
                    max_prompt_number = prompt_number

                if aggregated_data is None:
                    aggregated_data = df

                else:
                    aggregated_data = pd.merge(aggregated_data, df, on="iteration", how="outer")
            else:
                print(f"File not found: {file_path}")

    # Ensure aggregated_data is not None before saving
    if aggregated_data is not None:
        # Save the aggregated data to an Excel file
        output_file = os.path.join(OUTPUT_FOLDER, "aggregated_score_results.xlsx")
        aggregated_data.to_excel(output_file, index=False)
        print(f"Aggregated results saved to {output_file}")
    else:
        print("No data was aggregated. Check the input folders and files.")

    data = pd.read_excel(output_file)

    data['max_combined_score'] = aggregated_data[f"combined_score_{max_seed}_{max_prompt_number}"]
    data['max_combined_loss'] = aggregated_data[f"combined_loss_{max_seed}_{max_prompt_number}"]
    data['max_aesthetic_score'] = aggregated_data[f"aesthetic_score_{max_seed}_{max_prompt_number}"]
    data['max_clip_score'] = aggregated_data[f"clip_score_{max_seed}_{max_prompt_number}"]

    # Calculate the average metrics across all seeds for each iteration
    data['avg_combined_score'] = data.filter(like='combined_score_').mean(axis=1)
    data['avg_combined_loss'] = data.filter(like='combined_loss_').mean(axis=1)
    data['avg_aesthetic_score'] = data.filter(like='aesthetic_score_').mean(axis=1)
    data['avg_clip_score'] = data.filter(like='clip_score_').mean(axis=1)

    # Calculate the standard deviation for each metric
    data['std_combined_score'] = data.filter(like='combined_score_').std(axis=1)
    data['std_combined_loss'] = data.filter(like='combined_loss_').std(axis=1)
    data['std_aesthetic_score'] = data.filter(like='aesthetic_score_').std(axis=1)
    data['std_clip_score'] = data.filter(like='clip_score_').std(axis=1)

    # Loss Evolution
    plt.figure(figsize=(10, 6))
    plot_mean_std(data['iteration'], data['avg_combined_loss'], data['std_combined_loss'], "Loss")
    plt.plot(data['iteration'], data['max_combined_loss'], '-', label="Best")
    plt.ylim(0, 1.1)
    plt.xlabel('Iteration')
    plt.ylabel('Loss')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(OUTPUT_FOLDER + "/loss_evolution.png")
    plt.close()

    # Aesthetic Score Evolution
    plt.figure(figsize=(10, 6))
    plot_mean_std(data['iteration'], data['avg_aesthetic_score'], data['std_aesthetic_score'], "")
    plt.plot(data['iteration'], data['max_aesthetic_score'], '-', label="Best")
    plt.ylim(0, 10.5)
    plt.xlabel('Iteration')
    plt.ylabel('Aesthetic Score')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(OUTPUT_FOLDER + "/aesthetic_score_evolution.png")
    plt.close()

    # CLIP Score Evolution
    plt.figure()
    plot_mean_std(data['iteration'], data['avg_clip_score'], data['std_clip_score'], "Population")
    plt.plot(data['iteration'], data['max_clip_score'], '-', label="Best")
    plt.ylim(0, 0.6)
    plt.xlabel('Iteration')
    plt.ylabel('CLIP Score')
    plt.grid()
    plt.legend()
    plt.savefig(OUTPUT_FOLDER + "/clip_score_evolution.png")
    plt.close()

    # Initialize PowerPoint presentation
    presentation = Presentation()

    # Collect folders with seed numbers
    folders = []
    for folder_name in os.listdir(OUTPUT_FOLDER):
        if folder_name.startswith("results_"):
            # Extract the seed number from the folder name
            seed_number = int(folder_name.split("_")[-2])  # Convert to integer for sorting
            prompt_number = int(folder_name.split("_")[-1])  # Extract prompt number
            folder_path = os.path.join(OUTPUT_FOLDER, folder_name)
            folders.append((seed_number, prompt_number, folder_path))

    # Sort folders by prompt number in ascending order
    folders.sort(key=lambda x: x[1])

    # Iterate over sorted folders
    for seed_number, prompt_number, folder_path in folders:
        # Paths for required images and CSV file
        it_0_path = os.path.join(folder_path, "it_0.png")
        best_all_path = os.path.join(folder_path, "best_all.png")
        loss_evolution_path = os.path.join(folder_path, "loss_evolution.png")
        score_evolution_path = os.path.join(folder_path, "aesthetic_evolution.png")
        clip_evolution_path = os.path.join(folder_path, "clip_evolution.png")
        csv_path = os.path.join(folder_path, "score_results.csv")

        # Extract scores and prompt from CSV
        combined_score_initial = None
        combined_score_best = None
        aesthetic_initial = None
        aesthetic_best = None
        clip_initial = None
        clip_best = None
        prompt_text = None
        category = None

        if os.path.exists(csv_path):
            with open(csv_path, 'r') as csvfile:
                reader = csv.DictReader(csvfile)
                rows = list(reader)
                if rows:
                    # Initial values from the first row
                    first_row = rows[0]
                    combined_score_initial = float(first_row['combined_score'])
                    aesthetic_initial = float(first_row['aesthetic_score'])
                    clip_initial = float(first_row['clip_score'])
                    prompt_text = first_row['prompt']
                    category = first_row['category']

                    # Find the row with the best (maximum) max_fitness
                    best_row = max(rows, key=lambda r: float(r['combined_score']))
                    combined_score_best = float(best_row['combined_score'])
                    aesthetic_best = float(best_row['aesthetic_score'])
                    clip_best = float(best_row['clip_score'])

        # Slide 1: it_0.png and it_1000.png
        if os.path.exists(it_0_path) and os.path.exists(best_all_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            # Add prompt below the title
            if prompt_text:
                left = Inches(0.5)
                top = Inches(1)
                width = Inches(9)
                textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
                textbox.text = f"Prompt: {prompt_text}"

            if category:
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
                textbox.text = f"Category: {category}"

            # Add it_0.png
            slide.shapes.add_picture(it_0_path, Inches(0.5), Inches(2), height=Inches(4))

            # Add legend below it_0.png
            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Initial iteration"
            if combined_score_initial is not None:
                text += f"\nInitial Combined Score: {combined_score_initial:.4f}"
            if aesthetic_initial is not None:
                text += f"\nAesthetic Score: {aesthetic_initial:.4f}"
            if clip_initial is not None:
                text += f"\nCLIP Score: {clip_initial:.4f}"
            textbox.text = text

            # Add it_1000.png
            slide.shapes.add_picture(best_all_path, Inches(5.5), Inches(2), height=Inches(4))

            # Add legend below it_1000.png
            left = Inches(5.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Best iteration"
            if combined_score_best is not None:
                text += f"\nBest Combined Score: {combined_score_best:.4f}"
            if aesthetic_best is not None:
                text += f"\nAesthetic Score: {aesthetic_best:.4f}"
            if clip_best is not None:
                text += f"\nCLIP Score: {clip_best:.4f}"
            textbox.text = text

        # Slide 2: loss_evolution.png
        if os.path.exists(loss_evolution_path) and os.path.exists(score_evolution_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            # Add aesthetic_evolution.png
            slide.shapes.add_picture(loss_evolution_path, Inches(0), Inches(2), height=Inches(4))

            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Loss evolution"
            textbox.text = text

        # Slide 3: aesthetic_evolution.png and loss_evolution.png
        if os.path.exists(loss_evolution_path) and os.path.exists(score_evolution_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            # Add aesthetic_evolution.png
            slide.shapes.add_picture(clip_evolution_path, Inches(0), Inches(2), height=Inches(4))

            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "CLIP evolution"
            textbox.text = text

            # Add loss_evolution.png
            slide.shapes.add_picture(score_evolution_path, Inches(5), Inches(2), height=Inches(4))

            left = Inches(5.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Aesthetic evolution"
            textbox.text = text

    output_filename = os.path.join(OUTPUT_FOLDER, f"summary.pptx")
    # Save the presentation
    presentation.save(output_filename)
    print(f"Presentation saved as {output_filename}")

if __name__ == "__main__":
    seed_number = 1
    for seed in seed_list:
        prompt_number = 1
        for prompt, category in selected_prompts_with_category:
            print(f"Running seed {seed}, prompt: {prompt} (Category: {category})")
            main(seed, seed_number, prompt, category, prompt_number)
            print(f"Run with seed {seed} and prompt '{prompt}' finished!")
            aggregate_results()
            prompt_number += 1
        seed_number += 1
