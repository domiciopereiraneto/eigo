# promptist_baseline.py
import sys
import os
import shutil
import json
import yaml

# Get the parent directory
parent_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
sys.path.insert(0, parent_dir)

# External imports
import torch
import torch.nn.functional as F
import numpy as np
import pandas as pd
from diffusers import StableDiffusionXLPipeline
import random
from PIL import Image
import matplotlib.pyplot as plt
import time
from collections import defaultdict
import csv
from pptx import Presentation
from pptx.util import Inches
from deap import base, creator, tools
from datasets import load_dataset
import clip
import argparse
import warnings

from transformers import AutoModelForCausalLM, AutoTokenizer

# -----------------------
# Argument parsing
# -----------------------
parser = argparse.ArgumentParser(description='Run Promptist baseline with configuration file')
parser.add_argument('--config', type=str, default="algorithms/config/config_comparison.yaml",
                    help='Path to configuration YAML file')
args = parser.parse_args()

config_path = args.config

with open(config_path, 'r') as file:
    config = yaml.safe_load(file)

SEED = config['seed']
SEED_PATH = config['seed_path']
cuda_n = str(config['cuda'])
predictor = config['predictor']
num_inference_steps = config['num_inference_steps']
height = config['height']
width = config['width']
OUTPUT_FOLDER = config['results_folder']
alpha = config['alpha']
beta = config['beta']
max_aesthetic_score = config['max_aesthetic_score']
max_clip_score = config['max_clip_score']
model_id = config['model_id']

# Predictor name
if predictor == 0:
    predictor_name = 'simulacra'
elif predictor == 1:
    predictor_name = 'laionv1'
elif predictor == 2:
    predictor_name = 'laionv2'
else:
    raise ValueError("Invalid predictor option.")

# -----------------------
# Method name / output folder
# -----------------------
method_save_name = 'promptist'
OUTPUT_FOLDER = f"{OUTPUT_FOLDER}/{method_save_name}_clip_{predictor_name}_sdxlturbo_{SEED}_a{int(alpha*100)}_b{int(beta*100)}"

os.makedirs(OUTPUT_FOLDER, exist_ok=True)
shutil.copy(config_path, os.path.join(OUTPUT_FOLDER, "config_used.yaml"))

# -----------------------
# Device
# -----------------------
device = "cuda:" + cuda_n if torch.cuda.is_available() else "cpu"

# -----------------------
# Load SDXL(-Turbo) pipeline
# -----------------------
pipe = StableDiffusionXLPipeline.from_pretrained(
    model_id,
    torch_dtype=torch.float32,
    use_safetensors=True,
).to(device)
pipe.set_progress_bar_config(disable=True)

# -----------------------
# CLIP
# -----------------------
clip_model_name = "ViT-L/14"
clip_model, clip_preprocess = clip.load(clip_model_name, device=device)

# -----------------------
# Prompt dataset: Parti-prompts
# -----------------------
prompt_dataset = load_dataset("nateraw/parti-prompts")["train"]

N_PER_CATEGORY = config['prompt_per_categorie']
SUBSET_SEED = config['prompt_sample_seed']
random.seed(SUBSET_SEED)

category_prompts = defaultdict(list)
for item in prompt_dataset:
    category = item.get("Category", "Uncategorized")
    category_prompts[category].append(item["Prompt"])

selected_prompts_with_category = []
for category, prompts in category_prompts.items():
    if len(prompts) >= N_PER_CATEGORY:
        sampled = random.sample(prompts, N_PER_CATEGORY)
    else:
        sampled = prompts
    for prompt in sampled:
        selected_prompts_with_category.append((prompt, category))

prompt_list_path = os.path.join(OUTPUT_FOLDER, "selected_prompts.txt")
with open(prompt_list_path, "w", encoding="utf-8") as f:
    for prompt, category in selected_prompts_with_category:
        f.write(f"{category}\t{prompt}\n")
print(f"Saved selected prompts to {prompt_list_path}")
print(f"Selected {len(selected_prompts_with_category)} prompts from {len(category_prompts)} categories.")

# -----------------------
# Aesthetic model
# -----------------------
if predictor == 0:
    from aesthetic_evaluation.src import simulacra_rank_image
    aesthetic_model = simulacra_rank_image.SimulacraAesthetic(device)
    model_name = 'SAM'
elif predictor == 1:
    from aesthetic_evaluation.src import laion_rank_image
    aesthetic_model = laion_rank_image.LAIONAesthetic(device, clip_model=clip_model_name)
    model_name = 'LAIONV1'
elif predictor == 2:
    from aesthetic_evaluation.src import laion_v2_rank_image
    aesthetic_model = laion_v2_rank_image.LAIONV2Aesthetic(device, clip_model=clip_model_name)
    model_name = 'LAIONV2'
else:
    raise ValueError("Invalid predictor option.")

# -----------------------
# Seed handling
# -----------------------
if SEED_PATH is None:
    seed_list = [SEED]
else:
    with open(SEED_PATH, 'r') as file:
        seed_list = [int(line.strip()) for line in file]

# -----------------------
# Promptist
# -----------------------
print("Loading Promptist model...")
promptist_model = AutoModelForCausalLM.from_pretrained("microsoft/Promptist").to(device)
promptist_tokenizer = AutoTokenizer.from_pretrained("gpt2")
promptist_tokenizer.pad_token = promptist_tokenizer.eos_token
promptist_tokenizer.padding_side = "left"

@torch.no_grad()
def optimize_prompt(plain_text: str, top_k: int = 1):
    """
    Use Promptist to optimize a natural language prompt.
    """
    plain_text = plain_text.strip()
    input_ids = promptist_tokenizer(
        plain_text + " Rephrase:",
        return_tensors="pt"
    ).input_ids.to(device)

    eos_id = promptist_tokenizer.eos_token_id
    outputs = promptist_model.generate(
        input_ids,
        do_sample=False,
        max_new_tokens=75,
        num_beams=8,
        num_return_sequences=min(top_k, 8),
        eos_token_id=eos_id,
        pad_token_id=eos_id,
        length_penalty=-1.0,
    )

    texts = promptist_tokenizer.batch_decode(outputs, skip_special_tokens=True)
    cleaned = [
        t.replace(plain_text + " Rephrase:", "").strip()
        for t in texts
    ]
    if top_k == 1:
        return cleaned[0]
    return cleaned

# -----------------------
# Utility functions
# -----------------------
def pil_to_tensor(image: Image.Image) -> torch.Tensor:
    """
    Convert PIL image (H, W, 3) to float tensor in [0, 1], shape [H, W, C].
    """
    arr = np.array(image).astype(np.float32) / 255.0
    return torch.from_numpy(arr).to(device)

def generate_image_from_text(prompt: str, seed: int) -> Image.Image:
    """
    Generate an image from a text prompt using SDXL(-Turbo).
    """
    generator = torch.Generator(device=device).manual_seed(seed)
    with torch.no_grad():
        out = pipe(
            prompt=prompt,
            num_inference_steps=num_inference_steps,
            guidance_scale=0.0,
            height=height,
            width=width,
            generator=generator,
        ).images[0]
    return out

def aesthetic_evaluation(image_tensor: torch.Tensor) -> torch.Tensor:
    """
    Aesthetic score from tensor [H, W, C] in [0,1].
    """
    img_chw = image_tensor.permute(2, 0, 1).to(torch.float32)  # [C, H, W]

    if predictor in (0, 1, 2):
        score = aesthetic_model.predict_from_tensor(img_chw)
    else:
        score = torch.tensor(0.0, device=device)
    return score

def evaluate_clip_score(image_tensor: torch.Tensor, prompt: str) -> float:
    """
    CLIP score between image tensor [H, W, C] and text prompt.
    """
    image = (image_tensor * 255).clamp(0, 255).byte()
    image = Image.fromarray(image.detach().cpu().numpy())

    image_input = clip_preprocess(image).unsqueeze(0).to(device)
    text_input = clip.tokenize([prompt]).to(device)

    with torch.no_grad():
        image_features = clip_model.encode_image(image_input)
        text_features = clip_model.encode_text(text_input)

    image_features = image_features / image_features.norm(dim=-1, keepdim=True)
    text_features = text_features / text_features.norm(dim=-1, keepdim=True)

    clip_score = (image_features @ text_features.T).item()
    return clip_score

def format_time(seconds):
    seconds = int(seconds)
    hours = seconds // 3600
    minutes = (seconds % 3600) // 60
    seconds = seconds % 60
    if hours > 0:
        return f"{hours}h {minutes}m {seconds}s"
    elif minutes > 0:
        return f"{minutes}m {seconds}s"
    else:
        return f"{seconds}s"

def plot_mean_std(x_axis, m_vec, std_vec, description, title=None, y_label=None, x_label=None):
    lower_bound = [m - s for m, s in zip(m_vec, std_vec)]
    upper_bound = [m + s for m, s in zip(m_vec, std_vec)]

    plt.plot(x_axis, m_vec, '--', label=description + " Avg.")
    plt.fill_between(x_axis, lower_bound, upper_bound, alpha=.3, label=description + " Avg. ± SD")
    if title is not None:
        plt.title(title)
    if y_label is not None:
        plt.ylabel(y_label)
    if x_label is not None:
        plt.xlabel(x_label)

def save_plot_results(results, results_folder):
    plt.figure(figsize=(10, 6))
    plot_mean_std(results['generation'], results['avg_fitness'], results['std_fitness'], "Fitness")
    plt.plot(results['generation'], results['max_fitness'], 'r-', label="Best Fitness")
    plt.ylim(0, 1.1)
    plt.xlabel('Generation')
    plt.ylabel('Fitness')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(os.path.join(results_folder, "fitness_evolution.png"))
    plt.close()

    plt.figure(figsize=(10, 6))
    plot_mean_std(results['generation'], results['avg_aesthetic_score'], results['std_aesthetic_score'], "Population")
    plt.plot(results['generation'], results['max_aesthetic_score'], 'r-', label="Best")
    plt.ylim(0, 10)
    plt.xlabel('Generation')
    plt.ylabel('Aesthetic Score')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(os.path.join(results_folder, "aesthetic_score_evolution.png"))
    plt.close()

    plt.figure(figsize=(10, 6))
    plot_mean_std(results['generation'], results['avg_clip_score'], results['std_clip_score'], "Population")
    plt.plot(results['generation'], results['max_clip_score'], 'r-', label="Best")
    plt.ylim(0, 0.6)
    plt.xlabel('Generation')
    plt.ylabel('CLIP Score')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(os.path.join(results_folder, "clip_score_evolution.png"))
    plt.close()

def aggregate_results():
    aggregated_data = None

    for folder_name in os.listdir(OUTPUT_FOLDER):
        if folder_name.startswith(f"results_{model_name}_"):
            seed = folder_name.split("_")[-2]
            prompt_number = folder_name.split("_")[-1]

            file_path = os.path.join(OUTPUT_FOLDER, folder_name, "fitness_results.csv")
            if os.path.exists(file_path):
                df = pd.read_csv(file_path)

                df = df.drop(columns=["prompt", "category", "best_prompt"], errors='ignore')

                df = df.rename(columns={"avg_fitness": f"avg_fitness_{seed}_{prompt_number}"})
                df = df.rename(columns={"max_fitness": f"max_fitness_{seed}_{prompt_number}"})
                df = df.rename(columns={"std_fitness": f"std_fitness_{seed}_{prompt_number}"})
                df = df.rename(columns={"avg_aesthetic_score": f"avg_aesthetic_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"max_aesthetic_score": f"max_aesthetic_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"std_aesthetic_score": f"std_aesthetic_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"avg_clip_score": f"avg_clip_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"max_clip_score": f"max_clip_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"std_clip_score": f"std_clip_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"elapsed_time": f"elapsed_time_{seed}_{prompt_number}"})

                if aggregated_data is None:
                    aggregated_data = df
                else:
                    aggregated_data = pd.merge(aggregated_data, df, on="generation", how="outer")
            else:
                print(f"File not found: {file_path}")

    if aggregated_data is None:
        print("No data was aggregated. Check the input folders and files.")
        return

    output_file = os.path.join(OUTPUT_FOLDER, "aggregated_score_results.xlsx")
    aggregated_data.to_excel(output_file, index=False)
    print(f"Aggregated results saved to {output_file}")

    data = pd.read_excel(output_file)

    data['avg_fitness'] = data.filter(like='avg_fitness_').mean(axis=1)
    data['std_fitness'] = data.filter(like='std_fitness_').std(axis=1)
    data['best_avg_fitness'] = data.filter(like='max_fitness_').mean(axis=1)
    data['best_std_fitness'] = data.filter(like='max_fitness_').std(axis=1)
    data['max_fitness'] = data.filter(like='max_fitness_').max(axis=1)

    data['avg_aesthetic_score'] = data.filter(like='avg_aesthetic_score_').mean(axis=1)
    data['std_aesthetic_score'] = data.filter(like='std_aesthetic_score_').std(axis=1)
    data['best_avg_aesthetic_score'] = data.filter(like='max_aesthetic_score_').mean(axis=1)
    data['best_std_aesthetic_score'] = data.filter(like='max_aesthetic_score_').std(axis=1)
    data['max_aesthetic_score'] = data.filter(like='max_aesthetic_score_').max(axis=1)

    data['avg_clip_score'] = data.filter(like='avg_clip_score_').mean(axis=1)
    data['std_clip_score'] = data.filter(like='std_clip_score_').std(axis=1)
    data['best_avg_clip_score'] = data.filter(like='max_clip_score_').mean(axis=1)
    data['best_std_clip_score'] = data.filter(like='max_clip_score_').std(axis=1)
    data['max_clip_score'] = data.filter(like='max_clip_score_').max(axis=1)

    # Fitness evolution
    plt.figure(figsize=(10, 6))
    plot_mean_std(data['generation'], data['avg_fitness'], data['std_fitness'], "Population")
    plot_mean_std(data['generation'], data['best_avg_fitness'], data['best_std_fitness'], "Bests")
    plt.plot(data['generation'], data['max_fitness'], 'r-', label="Best")
    plt.ylim(0, 1.1)
    plt.xlabel('Generation')
    plt.ylabel('Fitness')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(os.path.join(OUTPUT_FOLDER, "fitness_evolution.png"))
    plt.close()

    # Aesthetic evolution
    plt.figure(figsize=(10, 6))
    plot_mean_std(data['generation'], data['avg_aesthetic_score'], data['std_aesthetic_score'], "Population")
    plot_mean_std(data['generation'], data['best_avg_aesthetic_score'], data['best_std_aesthetic_score'], "Bests")
    plt.plot(data['generation'], data['max_aesthetic_score'], 'r-', label="Best")
    plt.ylim(0, 10.5)
    plt.xlabel('Generation')
    plt.ylabel('Aesthetic Score')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(os.path.join(OUTPUT_FOLDER, "aesthetic_score_evolution.png"))
    plt.close()

    # CLIP evolution
    plt.figure(figsize=(10, 6))
    plot_mean_std(data['generation'], data['avg_clip_score'], data['std_clip_score'], "Population")
    plot_mean_std(data['generation'], data['best_avg_clip_score'], data['best_std_clip_score'], "Bests")
    plt.plot(data['generation'], data['max_clip_score'], 'r-', label="Best")
    plt.xlabel('Generation')
    plt.ylabel('CLIP Score')
    plt.ylim(0, 0.6)
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(os.path.join(OUTPUT_FOLDER, "clip_score_evolution.png"))
    plt.close()

    # PowerPoint summary (same style as GA script)
    presentation = Presentation()

    folders = []
    for folder_name in os.listdir(OUTPUT_FOLDER):
        if folder_name.startswith("results_"):
            seed_number = int(folder_name.split("_")[-2])
            prompt_number = int(folder_name.split("_")[-1])
            folder_path = os.path.join(OUTPUT_FOLDER, folder_name)
            folders.append((seed_number, prompt_number, folder_path))

    folders.sort(key=lambda x: x[1])

    for seed_number, prompt_number, folder_path in folders:
        it_0_path = os.path.join(folder_path, "it_0.png")
        best_all_path = os.path.join(folder_path, "best_all.png")
        fitness_evolution_path = os.path.join(folder_path, "fitness_evolution.png")
        aesthetic_evolution_path = os.path.join(folder_path, "aesthetic_score_evolution.png")
        clip_evolution_path = os.path.join(folder_path, "clip_score_evolution.png")
        csv_path = os.path.join(folder_path, "fitness_results.csv")

        fitness_initial = None
        fitness_best = None
        aesthetic_initial = None
        aesthetic_best = None
        clip_initial = None
        clip_best = None
        prompt_text = None
        category = None

        if os.path.exists(csv_path):
            with open(csv_path, 'r') as csvfile:
                reader = csv.DictReader(csvfile)
                rows = list(reader)
                if rows:
                    first_row = rows[0]
                    fitness_initial = float(first_row['max_fitness'])
                    aesthetic_initial = float(first_row['max_aesthetic_score'])
                    clip_initial = float(first_row['max_clip_score'])
                    prompt_text = first_row['prompt']
                    category = first_row['category']

                    best_row = max(rows, key=lambda r: float(r['max_fitness']))
                    fitness_best = float(best_row['max_fitness'])
                    aesthetic_best = float(best_row['max_aesthetic_score'])
                    clip_best = float(best_row['max_clip_score'])

        if os.path.exists(it_0_path) and os.path.exists(best_all_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            if prompt_text:
                left = Inches(0.5)
                top = Inches(1)
                width = Inches(9)
                textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
                textbox.text = f"Prompt: {prompt_text}"

            if category:
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
                textbox.text = f"Category: {category}"

            slide.shapes.add_picture(it_0_path, Inches(0.5), Inches(2), height=Inches(4))

            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Original prompt image"
            if fitness_initial is not None:
                text += f"\nFitness: {fitness_initial:.4f}"
            if aesthetic_initial is not None:
                text += f"\nAesthetic Score: {aesthetic_initial:.4f}"
            if clip_initial is not None:
                text += f"\nCLIP Score: {clip_initial:.4f}"
            textbox.text = text

            slide.shapes.add_picture(best_all_path, Inches(5.5), Inches(2), height=Inches(4))

            left = Inches(5.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Promptist prompt image"
            if fitness_best is not None:
                text += f"\nFitness: {fitness_best:.4f}"
            if aesthetic_best is not None:
                text += f"\nAesthetic Score: {aesthetic_best:.4f}"
            if clip_best is not None:
                text += f"\nCLIP Score: {clip_best:.4f}"
            textbox.text = text

        if os.path.exists(fitness_evolution_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"
            slide.shapes.add_picture(fitness_evolution_path, Inches(0), Inches(2), height=Inches(4))
            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            textbox.text = "Fitness evolution (orig vs Promptist)"

        if os.path.exists(clip_evolution_path) and os.path.exists(aesthetic_evolution_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            slide.shapes.add_picture(clip_evolution_path, Inches(0), Inches(2), height=Inches(4))
            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            textbox.text = "CLIP score evolution"

            slide.shapes.add_picture(aesthetic_evolution_path, Inches(5), Inches(2), height=Inches(4))
            left = Inches(5.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            textbox.text = "Aesthetic score evolution"

    output_filename = os.path.join(OUTPUT_FOLDER, "summary.pptx")
    presentation.save(output_filename)
    print(f"Presentation saved as {output_filename}")

# -----------------------
# Main per-prompt evaluation
# -----------------------
def main(seed, seed_number, selected_prompt, category, prompt_number):
    torch.manual_seed(seed)
    np.random.seed(seed)
    random.seed(seed)

    print(f"Selected prompt: {selected_prompt} (Category: {category})")

    results_folder = os.path.join(OUTPUT_FOLDER, f"results_{model_name}_{seed}_{prompt_number}")
    os.makedirs(results_folder, exist_ok=True)

    start_time = time.time()

    # ---- Generation 0: original prompt ----
    orig_image_pil = generate_image_from_text(selected_prompt, seed)
    orig_image_pil.save(os.path.join(results_folder, "it_0.png"))

    orig_tensor = pil_to_tensor(orig_image_pil)
    with torch.no_grad():
        orig_aesthetic = aesthetic_evaluation(orig_tensor).item()
        orig_clip = evaluate_clip_score(orig_tensor, selected_prompt)

    orig_fitness_1 = alpha * orig_aesthetic / max_aesthetic_score
    orig_fitness_2 = beta * orig_clip / max_clip_score
    orig_fitness = orig_fitness_1 + orig_fitness_2

    # ---- Generation 1: Promptist optimized prompt ----
    optimized_prompt = optimize_prompt(selected_prompt)
    print(f"Optimized prompt: {optimized_prompt}")

    promptist_image_pil = generate_image_from_text(optimized_prompt, seed)
    promptist_image_pil.save(os.path.join(results_folder, "best_all.png"))

    promptist_tensor = pil_to_tensor(promptist_image_pil)
    with torch.no_grad():
        opt_aesthetic = aesthetic_evaluation(promptist_tensor).item()
        opt_clip = evaluate_clip_score(promptist_tensor, optimized_prompt)

    opt_fitness_1 = alpha * opt_aesthetic / max_aesthetic_score
    opt_fitness_2 = beta * opt_clip / max_clip_score
    opt_fitness = opt_fitness_1 + opt_fitness_2

    elapsed_0 = 0.0
    elapsed_1 = time.time() - start_time

    generations = [0, 1]
    prompts_col = [selected_prompt, ""]
    category_col = [category, ""]

    avg_fitness_list = [orig_fitness, opt_fitness]
    std_fitness_list = [0.0, 0.0]
    max_fitness_list = [orig_fitness, opt_fitness]

    avg_aesthetic_list = [orig_aesthetic, opt_aesthetic]
    std_aesthetic_list = [0.0, 0.0]
    max_aesthetic_list = [orig_aesthetic, opt_aesthetic]

    avg_clip_list = [orig_clip, opt_clip]
    std_clip_list = [0.0, 0.0]
    max_clip_list = [orig_clip, opt_clip]

    best_prompt_list = [selected_prompt, optimized_prompt]
    time_list = [elapsed_0, elapsed_1]

    results = pd.DataFrame({
        "generation": generations,
        "prompt": prompts_col,
        "category": category_col,
        "avg_fitness": avg_fitness_list,
        "std_fitness": std_fitness_list,
        "max_fitness": max_fitness_list,
        "avg_aesthetic_score": avg_aesthetic_list,
        "std_aesthetic_score": std_aesthetic_list,
        "max_aesthetic_score": max_aesthetic_list,
        "avg_clip_score": avg_clip_list,
        "std_clip_score": std_clip_list,
        "max_clip_score": max_clip_list,
        "best_prompt": best_prompt_list,
        "elapsed_time": time_list
    })

    results.to_csv(os.path.join(results_folder, "fitness_results.csv"), index=False, na_rep='nan')
    save_plot_results(results, results_folder)

    print(
        f"Seed {seed_number} prompt #{prompt_number}: "
        f"orig_fitness={orig_fitness:.4f}, opt_fitness={opt_fitness:.4f}, "
        f"orig_aesthetic={orig_aesthetic:.4f}, opt_aesthetic={opt_aesthetic:.4f}, "
        f"orig_clip={orig_clip:.4f}, opt_clip={opt_clip:.4f}, "
        f"time={format_time(elapsed_1)}"
    )

# -----------------------
# Entry point
# -----------------------
if __name__ == "__main__":
    seed_number = 1
    for seed in seed_list:
        prompt_number = 1
        for prompt, category in selected_prompts_with_category:
            print(f"Running seed {seed}, prompt: {prompt} (Category: {category})")
            main(seed, seed_number, prompt, category, prompt_number)
            print(f"Run with seed {seed} and prompt '{prompt}' finished!")
            prompt_number += 1
        seed_number += 1

    aggregate_results()
