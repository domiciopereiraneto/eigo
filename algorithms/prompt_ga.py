# System imports
import sys
import os
import shutil
import json
import yaml

# Get the parent directory
parent_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
# Add the parent directory to sys.path to obtain access to the submodules
sys.path.insert(0, parent_dir)

# External imports - grouped by functionality
import torch
import torch.nn.functional as F
import numpy as np
import pandas as pd
from diffusers import StableDiffusionXLPipeline
import random
from PIL import Image
import matplotlib.pyplot as plt
import time
from collections import defaultdict
import csv
from pptx import Presentation
from pptx.util import Inches
from deap import base, creator, tools
from datasets import load_dataset
import clip
import argparse
import warnings

# Argument parsing for configuration file
# Allows specifying a custom configuration file path.
parser = argparse.ArgumentParser(description='Run optimization with configuration file')
parser.add_argument('--config', type=str, default="algorithms/config/config_ga.yaml",
                   help='Path to configuration YAML file')
args = parser.parse_args()

# Use the provided config path or default
config_path = args.config

# Load configuration parameters
# Reads the YAML configuration file and extracts parameters for the optimization process.
with open(config_path, 'r') as file:
    config = yaml.safe_load(file)

SEED = config['seed']
SEED_PATH = config['seed_path']
cuda_n = str(config['cuda'])
predictor = config['predictor']
num_inference_steps = config['num_inference_steps']
height = config['height']
width = config['width']
OUTPUT_FOLDER = config['results_folder']
NUM_GENERATIONS = config['num_generations']
POP_SIZE = config['pop_size']
IND_MUTATION_PROB = config['ind_mutation_prob']
TOURNMENT_SIZE = config['tournment_size']
ELITISM = config['elitism']
MUTATION_PROB = config['mutation_prob']
CROSSOVER_PROB = config['crossover_prob']
alpha = config['alpha']
beta = config['beta']
max_aesthetic_score = config['max_aesthetic_score']
max_clip_score = config['max_clip_score']
model_id = config['model_id']

# Determine the predictor and CMA-ES variant names based on the configuration
# Maps predictor indices and CMA-ES variants to their corresponding names.
if predictor == 0:
    predictor_name = 'simulacra'
elif predictor == 1:
    predictor_name = 'laionv1'
elif predictor == 2:
    predictor_name = 'laionv2'
else:
    raise ValueError("Invalid predictor option.")

method_save_name = 'ga'

# Set up the output folder
# Creates the output directory and saves the configuration file for reproducibility.
OUTPUT_FOLDER = f"{OUTPUT_FOLDER}/{method_save_name}_clip_{predictor_name}_sdxlturbo_{SEED}_a{int(alpha*100)}_b{int(beta*100)}"

# Save the selected prompts and their categories to a text file in the results folder
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
# Copy the YAML configuration file to the output folder
shutil.copy(config_path, os.path.join(OUTPUT_FOLDER, "config_used.yaml"))

# Check device availability
# Uses GPU if available, otherwise defaults to CPU.
device = "cuda:" + cuda_n if torch.cuda.is_available() else "cpu"

# Load the SDXL pipeline
# Initializes the pipeline for image generation with gradient computation enabled.
pipe = StableDiffusionXLPipeline.from_pretrained(
    model_id,
    torch_dtype=torch.float32,
    use_safetensors=True,
).to(device)
pipe.set_progress_bar_config(disable=True)

MIN_VALUE, MAX_VALUE = 0, pipe.tokenizer.vocab_size-3
START_OF_TEXT, END_OF_TEXT = pipe.tokenizer.bos_token_id, pipe.tokenizer.eos_token_id
VECTOR_SIZE = pipe.tokenizer.model_max_length

# DEAP setup for individuals and fitness; guard to avoid redefinition on repeated runs
if "FitnessMax" not in creator.__dict__:
    creator.create("FitnessMax", base.Fitness, weights=(1.0,))
if "Individual" not in creator.__dict__:
    creator.create("Individual", list, fitness=creator.FitnessMax)

# CLIP model setup
# Loads the CLIP model for evaluating image-text similarity.
clip_model_name = "ViT-L/14"  # CLIP model name
clip_model, clip_preprocess = clip.load(clip_model_name, device=device)

# Prompt dataset loading and preprocessing
# Groups prompts by category and samples a specified number per category.
prompt_dataset = load_dataset("nateraw/parti-prompts")["train"]

N_PER_CATEGORY = config['prompt_per_categorie']  # Number of prompts to sample per category
SUBSET_SEED = config['prompt_sample_seed']
random.seed(SUBSET_SEED)

# Group prompts by category
category_prompts = defaultdict(list)
for item in prompt_dataset:
    category = item.get("Category", "Uncategorized")
    category_prompts[category].append(item["Prompt"])

# Sample N_PER_CATEGORY prompts from each category and keep track of category
selected_prompts_with_category = []
for category, prompts in category_prompts.items():
    if len(prompts) >= N_PER_CATEGORY:
        sampled = random.sample(prompts, N_PER_CATEGORY)
    else:
        sampled = prompts  # If not enough, take all
    for prompt in sampled:
        selected_prompts_with_category.append((prompt, category))

# Save selected prompts to a file
# Stores the sampled prompts and their categories for reference.
prompt_list_path = os.path.join(OUTPUT_FOLDER, "selected_prompts.txt")
with open(prompt_list_path, "w", encoding="utf-8") as f:
    for prompt, category in selected_prompts_with_category:
        f.write(f"{category}\t{prompt}\n")
print(f"Saved selected prompts to {prompt_list_path}")

print(f"Selected {len(selected_prompts_with_category)} prompts from {len(category_prompts)} categories.")

# Initialize the aesthetic model
# Loads the appropriate aesthetic model based on the predictor configuration.
if predictor == 0:
    from aesthetic_evaluation.src import simulacra_rank_image
    aesthetic_model = simulacra_rank_image.SimulacraAesthetic(device)
    model_name = 'SAM'
elif predictor == 1:
    from aesthetic_evaluation.src import laion_rank_image
    aesthetic_model = laion_rank_image.LAIONAesthetic(device, clip_model=clip_model_name)
    model_name = 'LAIONV1'
elif predictor == 2:
    from aesthetic_evaluation.src import laion_v2_rank_image
    aesthetic_model = laion_v2_rank_image.LAIONV2Aesthetic(device, clip_model=clip_model_name)
    model_name = 'LAIONV2'
else:
    raise ValueError("Invalid predictor option.")

# Seed handling
# Initializes the random seed for reproducibility.
if SEED_PATH is None:
    seed_list = [SEED]
else:
    with open(SEED_PATH, 'r') as file:
        # Read each line, strip newline characters, and convert to integers
        seed_list = [int(line.strip()) for line in file]

def generate_image_from_prompt_tokens(token_vector, seed):
    generator = torch.Generator(device=device).manual_seed(seed)

    # Normalize token_vector → (1, seq_len) tensor of longs
    if isinstance(token_vector, dict):
        ids = token_vector["input_ids"]
    elif torch.is_tensor(token_vector):
        ids = token_vector
    else:
        ids = torch.tensor(token_vector, dtype=torch.long, device=device)

    ids = ids.view(1, -1).to(device)

    # Enforce fixed length and valid vocab range
    max_vocab = min(pipe.tokenizer.vocab_size, getattr(pipe, "tokenizer_2", pipe.tokenizer).vocab_size)
    ids = torch.clamp(ids[:, :VECTOR_SIZE], 0, max_vocab - 1)

    with torch.no_grad():
        # First text encoder
        enc_out_1 = pipe.text_encoder(ids, output_hidden_states=True)
        # Second text encoder (SDXL-style)
        enc_out_2 = pipe.text_encoder_2(ids, output_hidden_states=True)

        # SDXL uses the second-to-last hidden state as context for cross-attention
        emb_1 = enc_out_1.hidden_states[-2]  # (bs, seq, dim1)
        emb_2 = enc_out_2.hidden_states[-2]  # (bs, seq, dim2)
        prompt_embeds = torch.cat([emb_1, emb_2], dim=-1)  # (bs, seq, dim1+dim2)

        # Pooled embedding: SDXL normally uses the pooled output of the second encoder
        pooled_prompt_embeds = enc_out_2[0]  # (bs, pooled_dim)

    out = pipe(
        prompt_embeds=prompt_embeds,
        pooled_prompt_embeds=pooled_prompt_embeds,
        guidance_scale=0.0,
        num_inference_steps=num_inference_steps,
        generator=generator,
        output_type="pt",
    )["images"]

    image = out.clamp(0, 1).squeeze(0).permute(1, 2, 0)      # HWC
    return image.to(device)

def aesthetic_evaluation(image):
    """
    Evaluate the aesthetic quality of an image.

    Args:
        image (torch.Tensor): The image tensor of shape [H, W, C].

    Returns:
        torch.Tensor: The aesthetic score of the image.
    """
    # image is a tensor of shape [H, W, C]
    # Convert to [N, C, H, W] and ensure it's in float32
    image_input = image.permute(2, 0, 1).to(torch.float32)  # [1, C, H, W]

    if predictor == 0:
        # Simulacra Aesthetic Model
        score = aesthetic_model.predict_from_tensor(image_input)
    elif predictor == 1 or predictor == 2:
        # LAION Aesthetic Predictor V1 and V2
        score = aesthetic_model.predict_from_tensor(image_input)
    else:
        return torch.tensor(0.0, device=device)

    return score

def evaluate_clip_score(image_tensor, prompt):
    """
    Evaluate the alignment of a generated image with a text prompt using CLIP score.

    Args:
        image_tensor (torch.Tensor): The generated image as a tensor of shape [H, W, C] (values in [0, 1]).
        prompt (str): The text prompt to evaluate alignment with.
        device (str): The device to run the evaluation on ("cuda" or "cpu").

    Returns:
        float: The CLIP score indicating alignment between the image and the prompt.
    """

    # Convert the image tensor to a PIL image
    image = (image_tensor * 255).clamp(0, 255).byte()
    image = Image.fromarray(image.cpu().numpy())

    # Preprocess the image
    image_input = clip_preprocess(image).unsqueeze(0).to(device)

    # Tokenize the prompt
    text_input = clip.tokenize([prompt]).to(device)

    # Compute the CLIP embeddings
    image_features = clip_model.encode_image(image_input)
    text_features = clip_model.encode_text(text_input)

    # Normalize the features
    image_features = image_features / image_features.norm(dim=-1, keepdim=True)
    text_features = text_features / text_features.norm(dim=-1, keepdim=True)

    # Compute the cosine similarity (CLIP score)
    clip_score = (image_features @ text_features.T)

    return clip_score

def format_time(seconds):
    """
    Format the elapsed time in seconds to a human-readable string.

    Args:
        seconds (float): The elapsed time in seconds.

    Returns:
        str: The formatted time string.
    """
    seconds = int(seconds)
    hours = seconds // 3600
    minutes = (seconds % 3600) // 60
    seconds = seconds % 60
    if hours > 0:
        return f"{hours}h {minutes}m {seconds}s"
    elif minutes > 0:
        return f"{minutes}m {seconds}s"
    else:
        return f"{seconds}s"

def evaluate(token_vector, seed, selected_prompt, save_path=None):
    with torch.no_grad():
        image = generate_image_from_prompt_tokens(token_vector, seed)

        aesthetic_score = aesthetic_evaluation(image).item()
        clip_score = evaluate_clip_score(image, selected_prompt).item()
    # CMA-ES minimizes the function, so we need to invert the score if higher is better

    fitness_1 = alpha*aesthetic_score/max_aesthetic_score
    fitness_2 = beta*clip_score/max_clip_score

    fitness = fitness_1 + fitness_2

    if save_path is not None:
        # Save the generated image
        image_np = image.detach().clone().cpu().numpy()
        image_np = (image_np * 255).astype(np.uint8)
        pil_image = Image.fromarray(image_np)
        pil_image.save(save_path)

    return fitness, aesthetic_score, clip_score, fitness_1, fitness_2 

def detokenize(individual):
    tmp_solution = torch.tensor(individual, dtype=torch.int64)
    tmp_solution = torch.clamp(tmp_solution, 0, pipe.tokenizer.vocab_size - 1)
    decoded_string = pipe.tokenizer.decode(tmp_solution, skip_special_tokens=True, clean_up_tokenization_spaces = True)

    return decoded_string

def main(seed, seed_number, selected_prompt, category, prompt_number):
    """
    The main function for optimizing text embeddings and generating images.

    Args:
        seed (int): The random seed.
        seed_number (int): The seed number for tracking.
        selected_prompt (str): The selected text prompt.
        category (str): The category of the prompt.
        prompt_number (int): The prompt number for tracking.
    """
    torch.manual_seed(seed)
    np.random.seed(seed)
    random.seed(seed)

    print(f"Selected prompt: {selected_prompt} (Category: {category})")

    results_folder = f"{OUTPUT_FOLDER}/results_{model_name}_{seed}_{prompt_number}"
    os.makedirs(results_folder, exist_ok=True)

    # Get the initial token vector from the prompt
    initial_token_vector = pipe.tokenizer(
        selected_prompt,
        padding="max_length",
        max_length=pipe.tokenizer.model_max_length,
        truncation=True,
        return_tensors="pt",
    ).to(device)

    # Genetic Algorithm setup
    toolbox = base.Toolbox()

    # Register the mutation function
    toolbox.register("mutate", tools.mutUniformInt, low=MIN_VALUE, up=MAX_VALUE, indpb=IND_MUTATION_PROB)
    toolbox.register("mate", tools.cxOnePoint)
    toolbox.register("select", tools.selTournament, tournsize=TOURNMENT_SIZE)

    # Define a function to create an individual based on the initial token vector
    def create_individual_from_initial_vector(initial_vector):
        # Clone the initial vector to avoid modifying the original
        individual = initial_vector.clone().squeeze(0).tolist()
        # Apply mutation to the individual
        mutated_individual, = toolbox.mutate(individual)
        mutated_individual = mutated_individual[:VECTOR_SIZE]
        return creator.Individual(mutated_individual)

    # Register the population initialization function
    toolbox.register("individual", create_individual_from_initial_vector, initial_token_vector["input_ids"])
    toolbox.register("population", tools.initRepeat, list, toolbox.individual)
    toolbox.register("evaluate", evaluate, seed=seed, selected_prompt=selected_prompt)

    # Create the population
    population = toolbox.population(n=POP_SIZE)

    with torch.no_grad():
        initial_image = generate_image_from_prompt_tokens(initial_token_vector["input_ids"].squeeze(0), seed)
        image_np = initial_image.detach().clone().cpu().numpy()
        image_np = (image_np * 255).astype(np.uint8)
        pil_image = Image.fromarray(image_np)
        pil_image.save(f"{results_folder}/it_0.png")

        initial_fitness, initial_aesthetic_score, initial_clip_score, initial_fitness_1, initial_fitness_2 = evaluate(initial_token_vector["input_ids"].squeeze(0), seed, selected_prompt)

    time_list = [0]
    best_aesthetic_score_overall = initial_aesthetic_score
    best_clip_score_overall = initial_clip_score
    best_fitness_overall = initial_fitness
    best_tokens_overall = initial_token_vector

    start_time = time.time()
    generation = 0

    best_prompt_list = [selected_prompt]

    max_fit_list = [-initial_fitness]
    avg_fit_list = [-initial_fitness]
    std_fit_list = [0]

    max_aesthetic_score_list = [initial_aesthetic_score]
    avg_aesthetic_score_list = [initial_aesthetic_score]
    std_aesthetic_score_list = [0]

    max_clip_score_list = [initial_clip_score]
    avg_clip_score_list = [initial_clip_score]
    std_clip_score_list = [0]

    max_fitness_1_list = [initial_fitness_1]
    avg_fitness_1_list = [initial_fitness_1]
    std_fitness_1_list = [0]

    max_fitness_2_list = [initial_fitness_2]
    avg_fitness_2_list = [initial_fitness_2]
    std_fitness_2_list = [0]

    for gen in range(NUM_GENERATIONS):
        print(f"Generation {generation+1}/{NUM_GENERATIONS}")

        os.makedirs(results_folder+"/gen_%d" % (generation+1), exist_ok=True)

        # Evaluation metrics storage
        tmp_fitnesses = []
        tmp_fitness_1 = []
        tmp_fitness_2 = []
        aesthetic_scores = []
        clip_scores = []

        elites = tools.selBest(population, ELITISM)
        offspring = toolbox.select(population, len(population))
        offspring = list(map(toolbox.clone, offspring))

        for child1, child2 in zip(offspring[::2], offspring[1::2]):
            if random.random() < CROSSOVER_PROB:
                toolbox.mate(child1, child2)
                del child1.fitness.values
                del child2.fitness.values

        for mutant in offspring:
            if random.random() < MUTATION_PROB:
                toolbox.mutate(mutant)
                del mutant.fitness.values

        offspring = offspring + elites
        invalid_ind = [ind for ind in offspring if not ind.fitness.valid]
        fitnesses = list(map(toolbox.evaluate, invalid_ind))
        for ind, fit in zip(invalid_ind, fitnesses):
            ind.fitness.values = (fit[0],)
            ind.aesthetic_score = fit[1]
            ind.clip_score = fit[2]
            ind.fitness_1 = fit[3]
            ind.fitness_2 = fit[4]

        population[:] = offspring

        fits = [ind.fitness.values[0] for ind in population]
        max_fit = max(fits)
        avg_fit = sum(fits) / len(fits)
        std_fit = np.std(fits)

        aesthetic_scores = [ind.aesthetic_score for ind in population if hasattr(ind, "aesthetic_score")]
        clip_scores = [ind.clip_score for ind in population if hasattr(ind, "clip_score")]

        max_aesthetic_score_gen = max(aesthetic_scores)
        avg_aesthetic_score_gen = np.mean(aesthetic_scores)
        std_aesthetic_score_gen = np.std(aesthetic_scores)

        max_clip_score_gen = max(clip_scores)
        avg_clip_score_gen = np.mean(clip_scores)
        std_clip_score_gen = np.std(clip_scores)

        # Generate and display the best image
        best_ind = tools.selBest(population, 1)[0]
        prompt = detokenize(best_ind)

        best_prompt_list.append(prompt)
        max_fit_list.append(max_fit)
        avg_fit_list.append(avg_fit)
        std_fit_list.append(std_fit)
        max_aesthetic_score_list.append(max_aesthetic_score_gen)
        avg_aesthetic_score_list.append(avg_aesthetic_score_gen)
        std_aesthetic_score_list.append(std_aesthetic_score_gen)
        max_clip_score_list.append(max_clip_score_gen)
        avg_clip_score_list.append(avg_clip_score_gen)
        std_clip_score_list.append(std_clip_score_gen)

        # Get best solution so far
        best_x = prompt
        best_fitness = max_fit  # Convert back to positive score

        with torch.no_grad():
            # Generate and save the best image
            best_image = generate_image_from_prompt_tokens(best_ind, seed)
            image_np = best_image.detach().clone().cpu().numpy()
            image_np = (image_np * 255).astype(np.uint8)
            pil_image = Image.fromarray(image_np)
            pil_image.save(results_folder + "/best_%d.png" % (generation+1))

        if best_fitness > best_fitness_overall:
            best_fitness_overall = best_fitness
            best_text_embeddings_overall = best_x

        generation += 1

        elapsed_time = time.time() - start_time
        generations_done = generation
        generations_left = NUM_GENERATIONS - generations_done
        average_time_per_generation = elapsed_time / generations_done
        estimated_time_remaining = average_time_per_generation * generations_left

        formatted_time_remaining = format_time(estimated_time_remaining)

        time_list.append(elapsed_time)

        # Save the metrics
        results = pd.DataFrame({
            "generation": list(range(0, generation + 1)),
            "prompt": [selected_prompt] + [''] * generation,
            "category": [category] + [''] * generation,
            "avg_fitness": avg_fit_list,
            "std_fitness": std_fit_list,
            "max_fitness": max_fit_list,
            #"avg_fitness_1": avg_fitness_1_list,
            #"std_fitness_1": std_fitness_1_list,
            #"max_fitness_1": max_fitness_1_list,
            #"avg_fitness_2": avg_fitness_2_list,
            #"std_fitness_2": std_fitness_2_list,
            #"max_fitness_2": max_fitness_2_list,
            "avg_aesthetic_score": avg_aesthetic_score_list,
            "std_aesthetic_score": std_aesthetic_score_list,
            "max_aesthetic_score": max_aesthetic_score_list,
            "avg_clip_score": avg_clip_score_list,
            "std_clip_score": std_clip_score_list,
            "max_clip_score": max_clip_score_list,
            "best_prompt": best_prompt_list,
            "elapsed_time": time_list
        })

        results.to_csv(f"{results_folder}/fitness_results.csv", index=False, na_rep='nan')

        # Plot and save the fitness evolution
        save_plot_results(results, results_folder)

        #print(f"Gen {gen + 1}: Max fitness {max_fit}, Avg fitness {avg_fit}")
        # Print stats
        print(f"Seed {seed_number} Generation {generation}/{NUM_GENERATIONS}: Max fitness: {max_fit}, Avg fitness: {avg_fit}, Max aesthetic score: {max_aesthetic_score_gen}, Avg aesthetic score: {avg_aesthetic_score_list[-1]}, Max clip score: {max_clip_score_gen}, Avg clip score: {avg_clip_score_list[-1]}, Estimated time remaining: {formatted_time_remaining}")

    # Save the overall best image
    with torch.no_grad():
        best_image = generate_image_from_prompt_tokens(best_ind, seed)
    best_image_np = best_image.detach().cpu().numpy()
    best_image_np = (best_image_np * 255).astype(np.uint8)
    pil_image = Image.fromarray(best_image_np)
    pil_image.save(f"{results_folder}/best_all.png")

def plot_mean_std(x_axis, m_vec, std_vec, description, title=None, y_label=None, x_label=None):
    """
    Plot the mean and standard deviation with optional labels and title.

    Args:
        x_axis (iterable): The x-axis values.
        m_vec (iterable): The mean values.
        std_vec (iterable): The standard deviation values.
        description (str): A description for the plot legend.
        title (str): The title of the plot (optional).
        y_label (str): The label for the y-axis (optional).
        x_label (str): The label for the x-axis (optional).
    """
    lower_bound = [M_new - Sigma for M_new, Sigma in zip(m_vec, std_vec)]
    upper_bound = [M_new + Sigma for M_new, Sigma in zip(m_vec, std_vec)]

    plt.plot(x_axis, m_vec, '--', label=description + " Avg.")
    plt.fill_between(x_axis, lower_bound, upper_bound, alpha=.3, label=description + " Avg. ± SD")
    if title is not None:
        plt.title(title)
    if y_label is not None:
        plt.ylabel(y_label)
    if x_label is not None:
        plt.xlabel(x_label)

def save_plot_results(results, results_folder):
    """
    Generate and save plots for the evolution of scores and losses over generations.

    Args:
        results (pd.DataFrame): The DataFrame containing the results data.
        results_folder (str): The folder path to save the plots.
    """
    # Plot main fitness evolution
    plt.figure(figsize=(10, 6))  # Increase figure size
    plot_mean_std(results['generation'], results['avg_fitness'], results['std_fitness'], "Fitness")
    plt.plot(results['generation'], results['max_fitness'], 'r-', label="Best Fitness")
    #plot_mean_std(results['generation'], results['avg_fitness_1'], results['std_fitness_1'], "F1 (Aesthetic Score)")
    #plt.plot(results['generation'], results['max_fitness_1'], 'orange', label="Best F1")
    #plot_mean_std(results['generation'], results['avg_fitness_2'], results['std_fitness_2'], "F2 (CLIP Score)")
    #plt.plot(results['generation'], results['max_fitness_2'], 'green', label="Best F2")
    plt.ylim(0, 1.1)
    plt.xlabel('Generation')
    plt.ylabel('Fitness')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))  # Move legend outside the plot
    plt.tight_layout()  # Adjust layout
    plt.savefig(results_folder + "/fitness_evolution.png")
    plt.close()

    # Plot aesthetic score evolution
    plt.figure(figsize=(10, 6))  # Increase figure size
    plot_mean_std(results['generation'], results['avg_aesthetic_score'], results['std_aesthetic_score'], "Population")
    plt.plot(results['generation'], results['max_aesthetic_score'], 'r-', label="Best")
    plt.ylim(0, 10)
    plt.xlabel('Generation')
    plt.ylabel('Aesthetic Score')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))  # Move legend outside the plot
    plt.tight_layout()  # Adjust layout
    plt.savefig(results_folder + "/aesthetic_score_evolution.png")
    plt.close()

    # Plot clip score evolution
    plt.figure(figsize=(10, 6))  # Increase figure size
    plot_mean_std(results['generation'], results['avg_clip_score'], results['std_clip_score'], "Population")
    plt.plot(results['generation'], results['max_clip_score'], 'r-', label="Best")
    plt.ylim(0, 0.6)
    plt.xlabel('Generation')
    plt.ylabel('CLIP Score')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))  # Move legend outside the plot
    plt.tight_layout()  # Adjust layout
    plt.savefig(results_folder + "/clip_score_evolution.png")
    plt.close()

def aggregate_results():
    """
    Combine results from multiple runs and calculate summary statistics.

    Aggregates fitness, aesthetic, and CLIP scores across different seeds and prompts.
    Saves the aggregated results to an Excel file and generates summary plots.
    """
    # Initialize aggregated_data as None
    aggregated_data = None

    # Iterate over all subdirectories
    for folder_name in os.listdir(OUTPUT_FOLDER):
        if folder_name.startswith(f"results_{model_name}_"):
            seed = folder_name.split("_")[-2]  # Extract the seed number
            prompt_number = folder_name.split("_")[-1]  # Extract the prompt number

            file_path = os.path.join(OUTPUT_FOLDER, folder_name, "fitness_results.csv")

            if os.path.exists(file_path):
                # Read the CSV file
                df = pd.read_csv(file_path)

                df = df.drop(columns=["prompt", "category", "best_prompt"]) 

                df = df.rename(columns={"avg_fitness": f"avg_fitness_{seed}_{prompt_number}"})
                df = df.rename(columns={"max_fitness": f"max_fitness_{seed}_{prompt_number}"})
                df = df.rename(columns={"std_fitness": f"std_fitness_{seed}_{prompt_number}"})
                #df = df.rename(columns={"avg_fitness_1": f"avg_fitness_1_{seed}_{prompt_number}"})
                #df = df.rename(columns={"max_fitness_1": f"max_fitness_1_{seed}_{prompt_number}"})
                #df = df.rename(columns={"std_fitness_1": f"std_fitness_1_{seed}_{prompt_number}"})
                #df = df.rename(columns={"avg_fitness_2": f"avg_fitness_2_{seed}_{prompt_number}"})
                #df = df.rename(columns={"max_fitness_2": f"max_fitness_2_{seed}_{prompt_number}"})
                #df = df.rename(columns={"std_fitness_2": f"std_fitness_2_{seed}_{prompt_number}"})
                df = df.rename(columns={"avg_aesthetic_score": f"avg_aesthetic_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"max_aesthetic_score": f"max_aesthetic_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"std_aesthetic_score": f"std_aesthetic_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"avg_clip_score": f"avg_clip_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"max_clip_score": f"max_clip_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"std_clip_score": f"std_clip_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"elapsed_time": f"elapsed_time_{seed}_{prompt_number}"})

                if aggregated_data is None:
                    aggregated_data = df

                aggregated_data = pd.merge(aggregated_data, df, on="generation", how="outer")
            else:
                print(f"File not found: {file_path}")

    # Ensure aggregated_data is not None before saving
    if aggregated_data is not None:
        # Save the aggregated data to an Excel file
        output_file = os.path.join(OUTPUT_FOLDER, "aggregated_score_results.xlsx")
        aggregated_data.to_excel(output_file, index=False)
        print(f"Aggregated results saved to {output_file}")
    else:
        print("No data was aggregated. Check the input folders and files.")

    data = pd.read_excel(output_file)

    # Calculate the average fitness across all seeds for each iteration
    data['avg_fitness'] = data.filter(like='avg_fitness_').mean(axis=1)
    # Calculate the standard deviation of fitness across all seeds for each iteration
    data['std_fitness'] = data.filter(like='std_fitness_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_avg_fitness'] = data.filter(like='max_fitness_').mean(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_std_fitness'] = data.filter(like='max_fitness_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['max_fitness'] = data.filter(like='max_fitness_').max(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['avg_aesthetic_score'] = data.filter(like='avg_aesthetic_score_').mean(axis=1)
    # Calculate the standard deviation of fitness across all seeds for each iteration
    data['std_aesthetic_score'] = data.filter(like='std_aesthetic_score_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_avg_aesthetic_score'] = data.filter(like='max_aesthetic_score_').mean(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_std_aesthetic_score'] = data.filter(like='max_aesthetic_score_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['max_aesthetic_score'] = data.filter(like='max_aesthetic_score_').max(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['avg_clip_score'] = data.filter(like='avg_clip_score_').mean(axis=1)
    # Calculate the standard deviation of fitness across all seeds for each iteration
    data['std_clip_score'] = data.filter(like='std_clip_score_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_avg_clip_score'] = data.filter(like='max_clip_score_').mean(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_std_clip_score'] = data.filter(like='max_clip_score_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['max_clip_score'] = data.filter(like='max_clip_score_').max(axis=1)

    # Fitness Evolution
    plt.figure(figsize=(10, 6))
    plot_mean_std(data['generation'], data['avg_fitness'], data['std_fitness'], "Population")
    plot_mean_std(data['generation'], data['best_avg_fitness'], data['best_std_fitness'], "Bests")
    plt.plot(data['generation'], data['max_fitness'], 'r-', label="Best")
    plt.ylim(0, 1.1)
    plt.xlabel('Generation')
    plt.ylabel('Fitness')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(OUTPUT_FOLDER + "/fitness_evolution.png")
    plt.close()

    # Aesthetic Score Evolution
    plt.figure(figsize=(10, 6))
    plot_mean_std(data['generation'], data['avg_aesthetic_score'], data['std_aesthetic_score'], "Population")
    plot_mean_std(data['generation'], data['best_avg_aesthetic_score'], data['best_std_aesthetic_score'], "Bests")
    plt.plot(data['generation'], data['max_aesthetic_score'], 'r-', label="Best")
    plt.ylim(0, 10.5)
    plt.xlabel('Generation')
    plt.ylabel('Aesthetic Score')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(OUTPUT_FOLDER + "/aesthetic_score_evolution.png")
    plt.close()

    # CLIP Score Evolution
    plt.figure(figsize=(10, 6))
    plot_mean_std(data['generation'], data['avg_clip_score'], data['std_clip_score'], "Population")
    plot_mean_std(data['generation'], data['best_avg_clip_score'], data['best_std_clip_score'], "Bests")
    plt.plot(data['generation'], data['max_clip_score'], 'r-', label="Best")
    plt.xlabel('Generation')
    plt.ylabel('CLIP Score')
    plt.ylim(0, 0.6)
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(OUTPUT_FOLDER + "/clip_score_evolution.png")
    plt.close()

    # Initialize PowerPoint presentation
    presentation = Presentation()

    # Collect folders with seed numbers
    folders = []
    for folder_name in os.listdir(OUTPUT_FOLDER):
        if folder_name.startswith("results_"):
            # Extract the seed number from the folder name
            seed_number = int(folder_name.split("_")[-2])  # Convert to integer for sorting
            prompt_number = int(folder_name.split("_")[-1])  # Extract prompt number
            folder_path = os.path.join(OUTPUT_FOLDER, folder_name)
            folders.append((seed_number, prompt_number, folder_path))

    # Sort folders by prompt number in ascending order
    folders.sort(key=lambda x: x[1])

    # Iterate over sorted folders
    for seed_number, prompt_number, folder_path in folders:
        # Paths for required images and CSV file
        it_0_path = os.path.join(folder_path, "it_0.png")
        best_all_path = os.path.join(folder_path, "best_all.png")
        fitness_evolution_path = os.path.join(folder_path, "fitness_evolution.png")
        aesthetic_evolution_path = os.path.join(folder_path, "aesthetic_score_evolution.png")
        clip_evolution_path = os.path.join(folder_path, "clip_score_evolution.png")
        csv_path = os.path.join(folder_path, "fitness_results.csv")

        # Extract scores and prompt from CSV
        fitness_initial = None
        fitness_best = None
        aesthetic_initial = None
        aesthetic_best = None
        clip_initial = None
        clip_best = None
        prompt_text = None
        category = None

        if os.path.exists(csv_path):
            with open(csv_path, 'r') as csvfile:
                reader = csv.DictReader(csvfile)
                rows = list(reader)
                if rows:
                    # Initial values from the first row
                    first_row = rows[0]
                    fitness_initial = float(first_row['max_fitness'])
                    aesthetic_initial = float(first_row['max_aesthetic_score'])
                    clip_initial = float(first_row['max_clip_score'])
                    prompt_text = first_row['prompt']
                    category = first_row['category']

                    # Find the row with the best (maximum) max_fitness
                    best_row = max(rows, key=lambda r: float(r['max_fitness']))
                    fitness_best = float(best_row['max_fitness'])
                    aesthetic_best = float(best_row['max_aesthetic_score'])
                    clip_best = float(best_row['max_clip_score'])

        # Slide 1: it_0.png and it_1000.png
        if os.path.exists(it_0_path) and os.path.exists(best_all_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            # Add prompt below the title
            if prompt_text:
                left = Inches(0.5)
                top = Inches(1)
                width = Inches(9)
                textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
                textbox.text = f"Prompt: {prompt_text}"

            if category:
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
                textbox.text = f"Category: {category}"

            # Add it_0.png
            slide.shapes.add_picture(it_0_path, Inches(0.5), Inches(2), height=Inches(4))

            # Add legend below it_0.png
            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Initial iteration"
            if fitness_initial is not None:
                text += f"\nInitial Fitness: {fitness_initial:.4f}"
            if aesthetic_initial is not None:
                text += f"\nAesthetic Score: {aesthetic_initial:.4f}"
            if clip_initial is not None:
                text += f"\nCLIP Score: {clip_initial:.4f}"
            textbox.text = text

            # Add it_1000.png
            slide.shapes.add_picture(best_all_path, Inches(5.5), Inches(2), height=Inches(4))

            # Add legend below it_1000.png
            left = Inches(5.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Best iteration"
            if fitness_best is not None:
                text += f"\nBest fitness: {fitness_best:.4f}"
            if aesthetic_best is not None:
                text += f"\nAesthetic Score: {aesthetic_best:.4f}"
            if clip_best is not None:
                text += f"\nCLIP Score: {clip_best:.4f}"
            textbox.text = text

        # Slide 2: fitness_evolution
        if os.path.exists(fitness_evolution_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            # Add aesthetic_evolution.png
            slide.shapes.add_picture(fitness_evolution_path, Inches(0), Inches(2), height=Inches(4))

            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Fitness evolution"
            textbox.text = text

        # Slide 3: clip_score_evolution.png and aesthetic_score_evolution.png
        if os.path.exists(clip_evolution_path) and os.path.exists(aesthetic_evolution_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            # Add aesthetic_evolution.png
            slide.shapes.add_picture(clip_evolution_path, Inches(0), Inches(2), height=Inches(4))

            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "CLIP score evolution"
            textbox.text = text

            # Add loss_evolution.png
            slide.shapes.add_picture(aesthetic_evolution_path, Inches(5), Inches(2), height=Inches(4))

            left = Inches(5.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Aesthetic score evolution"
            textbox.text = text

    output_filename = os.path.join(OUTPUT_FOLDER, f"summary.pptx")
    # Save the presentation
    presentation.save(output_filename)
    print(f"Presentation saved as {output_filename}")

if __name__ == "__main__":
    # Entry point for the script
    # Parses arguments, loads configuration, and starts the optimization process.
    seed_number = 1
    for seed in seed_list:
        prompt_number = 1
        for prompt, category in selected_prompts_with_category:
            print(f"Running seed {seed}, prompt: {prompt} (Category: {category})")
            main(seed, seed_number, prompt, category, prompt_number)
            print(f"Run with seed {seed} and prompt '{prompt}' finished!")
            aggregate_results()
            prompt_number += 1
        seed_number += 1
