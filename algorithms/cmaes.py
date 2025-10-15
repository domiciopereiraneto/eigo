"""
CMA-ES-based optimization of text embeddings for image generation using SDXL.
Uses CMA-ES (including standard CMA-ES, sep-CMA-ES, or VD-CMA) optimizer to modify 
text embeddings while maximizing aesthetic and CLIP scores.
"""

import sys
import os
import shutil
import json
import yaml

# Get the parent directory
parent_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
# Add the parent directory to sys.path to obtain access to the submodules
sys.path.insert(0, parent_dir)

import torch
import torch.nn.functional as F
import numpy as np
import pandas as pd
from diffusers import StableDiffusionXLPipeline
import random
from PIL import Image
import matplotlib.pyplot as plt
import time
from collections import defaultdict
import csv
from pptx import Presentation
from pptx.util import Inches
import cma
from cma.restricted_gaussian_sampler import GaussVDSampler 
from datasets import load_dataset
import clip
import argparse

# Add argument parsing
parser = argparse.ArgumentParser(description='Run optimization with configuration file')
parser.add_argument('--config', type=str, default="algorithms/config/config_cmaes.yaml",
                   help='Path to configuration YAML file')
args = parser.parse_args()

# Use the provided config path or default
config_path = args.config

# Load configuration from YAML file
with open(config_path, 'r') as file:
    config = yaml.safe_load(file)

SEED = config['seed']
SEED_PATH = config['seed_path']
cuda_n = str(config['cuda'])
predictor = config['predictor']
num_inference_steps = config['num_inference_steps']
height = config['height']
width = config['width']
OUTPUT_FOLDER = config['results_folder']
NUM_GENERATIONS = config['num_generations']
POP_SIZE = config['pop_size']
SIGMA = config['sigma']
alpha = config['alpha']
beta = config['beta']
max_aesthetic_score = config['max_aesthetic_score']
max_clip_score = config['max_clip_score']
model_id = config['model_id']
cmaes_variant = config['cmaes_variant']

if predictor == 0:
    predictor_name = 'simulacra'
elif predictor == 1:
    predictor_name = 'laionv1'
elif predictor == 2:
    predictor_name = 'laionv2'
else:
    raise ValueError("Invalid predictor option.")

if cmaes_variant == "cmaes":
    method_save_name = "cmaes"
elif cmaes_variant == "sep":
    method_save_name = "sepcmaes"
elif cmaes_variant == "vd":
    method_save_name = "vdcmae"
else:
    raise ValueError(f"Unknown CMA-ES variant: {cmaes_variant}")

OUTPUT_FOLDER = f"{OUTPUT_FOLDER}/{method_save_name}_clip_{predictor_name}_sdxlturbo_{SEED}_a{int(alpha*100)}_b{int(beta*100)}"

# Save the selected prompts and their categories to a text file in the results folder
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
# Copy the YAML configuration file to the output folder
shutil.copy(config_path, os.path.join(OUTPUT_FOLDER, "config_used.yaml"))

# Check if a GPU is available and if not, use the CPU
device = "cuda:" + cuda_n if torch.cuda.is_available() else "cpu"

# Load the SDXL pipeline
pipe = StableDiffusionXLPipeline.from_pretrained(
    model_id,
    torch_dtype=torch.float32,
    use_safetensors=True,
).to(device)
pipe.set_progress_bar_config(disable=True)

clip_model_name = "ViT-L/14"  # CLIP model name
clip_model, clip_preprocess = clip.load(clip_model_name, device=device)

prompt_dataset = load_dataset("nateraw/parti-prompts")["train"]

N_PER_CATEGORY = config['prompt_per_categorie']  # Number of prompts to sample per category
SUBSET_SEED = config['prompt_sample_seed']
random.seed(SUBSET_SEED)

# Group prompts by category
category_prompts = defaultdict(list)
for item in prompt_dataset:
    category = item.get("Category", "Uncategorized")
    category_prompts[category].append(item["Prompt"])

# Sample N_PER_CATEGORY prompts from each category and keep track of category
selected_prompts_with_category = []
for category, prompts in category_prompts.items():
    if len(prompts) >= N_PER_CATEGORY:
        sampled = random.sample(prompts, N_PER_CATEGORY)
    else:
        sampled = prompts  # If not enough, take all
    for prompt in sampled:
        selected_prompts_with_category.append((prompt, category))

# Save the selected prompts and their categories to a text file in the results folder
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
prompt_list_path = os.path.join(OUTPUT_FOLDER, "selected_prompts.txt")
with open(prompt_list_path, "w", encoding="utf-8") as f:
    for prompt, category in selected_prompts_with_category:
        f.write(f"{category}\t{prompt}\n")
print(f"Saved selected prompts to {prompt_list_path}")

print(f"Selected {len(selected_prompts_with_category)} prompts from {len(category_prompts)} categories.")

# Initialize the aesthetic model
if predictor == 0:
    from aesthetic_evaluation.src import simulacra_rank_image
    aesthetic_model = simulacra_rank_image.SimulacraAesthetic(device)
    model_name = 'SAM'
elif predictor == 1:
    from aesthetic_evaluation.src import laion_rank_image
    aesthetic_model = laion_rank_image.LAIONAesthetic(device, clip_model=clip_model_name)
    model_name = 'LAIONV1'
elif predictor == 2:
    from aesthetic_evaluation.src import laion_v2_rank_image
    aesthetic_model = laion_v2_rank_image.LAIONV2Aesthetic(device, clip_model=clip_model_name)
    model_name = 'LAIONV2'
else:
    raise ValueError("Invalid predictor option.")

if SEED_PATH is None:
    seed_list = [SEED]
else:
    with open(SEED_PATH, 'r') as file:
        # Read each line, strip newline characters, and convert to integers
        seed_list = [int(line.strip()) for line in file]

def generate_image_from_embeddings(prompt_embeds, pooled_prompt_embeds, seed):
    generator = torch.Generator(device=device).manual_seed(seed)

    out = pipe(
        prompt_embeds              = prompt_embeds,
        pooled_prompt_embeds       = pooled_prompt_embeds,
        guidance_scale             = 0.0,
        num_inference_steps        = num_inference_steps,
        generator                  = generator,
        height                     = height,
        width                      = width,
        output_type               = "pt"
    )["images"]

    image = out.clamp(0, 1).squeeze(0).permute(1, 2, 0)      # HWC
    return image.to(device)

def aesthetic_evaluation(image):
    # image is a tensor of shape [H, W, C]
    # Convert to [N, C, H, W] and ensure it's in float32
    image_input = image.permute(2, 0, 1).to(torch.float32)  # [1, C, H, W]

    if predictor == 0:
        # Simulacra Aesthetic Model
        score = aesthetic_model.predict_from_tensor(image_input)
    elif predictor == 1 or predictor == 2:
        # LAION Aesthetic Predictor V1 and V2
        score = aesthetic_model.predict_from_tensor(image_input)
    else:
        return torch.tensor(0.0, device=device)

    return score

def evaluate_clip_score(image_tensor, prompt):
    """
    Evaluate the alignment of a generated image with a text prompt using CLIP score.

    Args:
        image_tensor (torch.Tensor): The generated image as a tensor of shape [H, W, C] (values in [0, 1]).
        prompt (str): The text prompt to evaluate alignment with.
        device (str): The device to run the evaluation on ("cuda" or "cpu").

    Returns:
        float: The CLIP score indicating alignment between the image and the prompt.
    """

    # Convert the image tensor to a PIL image
    image = (image_tensor * 255).clamp(0, 255).byte()
    image = Image.fromarray(image.cpu().numpy())

    # Preprocess the image
    image_input = clip_preprocess(image).unsqueeze(0).to(device)

    # Tokenize the prompt
    text_input = clip.tokenize([prompt]).to(device)

    # Compute the CLIP embeddings
    image_features = clip_model.encode_image(image_input)
    text_features = clip_model.encode_text(text_input)

    # Normalize the features
    image_features = image_features / image_features.norm(dim=-1, keepdim=True)
    text_features = text_features / text_features.norm(dim=-1, keepdim=True)

    # Compute the cosine similarity (CLIP score)
    clip_score = (image_features @ text_features.T)

    return clip_score

def format_time(seconds):
    seconds = int(seconds)
    hours = seconds // 3600
    minutes = (seconds % 3600) // 60
    seconds = seconds % 60
    if hours > 0:
        return f"{hours}h {minutes}m {seconds}s"
    elif minutes > 0:
        return f"{minutes}m {seconds}s"
    else:
        return f"{seconds}s"

def evaluate(input_embedding, seed, embedding_shape, selected_prompt, save_path=None):
    # x is a NumPy array representing the embedding vector
    # Convert it to a torch tensor

    # Reshape the embedding to the original shape
    split = np.prod(embedding_shape[0])
    pe  = torch.tensor(input_embedding[:split],  dtype=torch.float32, device=device).view(embedding_shape[0])
    ppe = torch.tensor(input_embedding[split:], dtype=torch.float32, device=device).view(embedding_shape[1])

    with torch.no_grad():
        image = generate_image_from_embeddings(pe, ppe, seed)

        aesthetic_score = aesthetic_evaluation(image).item()
        clip_score = evaluate_clip_score(image, selected_prompt).item()
    # CMA-ES minimizes the function, so we need to invert the score if higher is better

    fitness_1 = alpha*aesthetic_score/max_aesthetic_score
    fitness_2 = beta*clip_score/max_clip_score

    fitness = fitness_1 + fitness_2

    if save_path is not None:
        # Save the generated image
        image_np = image.detach().clone().cpu().numpy()
        image_np = (image_np * 255).astype(np.uint8)
        pil_image = Image.fromarray(image_np)
        pil_image.save(save_path)

    return -fitness, aesthetic_score, clip_score, fitness_1, fitness_2 

def main(seed, seed_number, selected_prompt, category, prompt_number):
    torch.manual_seed(seed)
    np.random.seed(seed)
    random.seed(seed)

    print(f"Selected prompt: {selected_prompt} (Category: {category})")

    results_folder = f"{OUTPUT_FOLDER}/results_{model_name}_{seed}_{prompt_number}"
    os.makedirs(results_folder, exist_ok=True)

    with torch.no_grad():
        (prompt_embeds,
        negative_prompt_embeds,
        pooled_prompt_embeds,
        negative_pooled_prompt_embeds) = pipe.encode_prompt(
                selected_prompt,
                negative_prompt="",                 
                device=device,
                num_images_per_prompt=1,
                do_classifier_free_guidance=False
        )

    # Set CMA-ES options
    es_options = {
        'seed': seed,
        'popsize': POP_SIZE,
        'maxiter': NUM_GENERATIONS,
        'verb_filenameprefix': results_folder + '/outcmaes',  # Save logs
        'verb_log': 0,  # Disable log output
        'verbose': -9,  # Suppress console output
    }

    if cmaes_variant == "cmaes":
        print("Using standard CMA-ES")
    elif cmaes_variant == "sep":
        print("Using sep-CMA-ES")
        es_options['CMA_diagonal'] = True
    elif cmaes_variant == "vd":
        print("Using VD-CMA-ES")
        es_options = GaussVDSampler.extend_cma_options(es_options)
    else:
        raise ValueError(f"Unknown CMA-ES variant: {cmaes_variant}")

    trainable_params_init = torch.cat([
        prompt_embeds.flatten(),
        pooled_prompt_embeds.flatten()
        ]).cpu().numpy()

    sh_prompt  = prompt_embeds.shape         
    sh_pooled  = pooled_prompt_embeds.shape  

    text_embeddings_init_shape = [sh_prompt, sh_pooled]

    es = cma.CMAEvolutionStrategy(trainable_params_init, SIGMA, es_options)

    with torch.no_grad():
        initial_image = generate_image_from_embeddings(prompt_embeds.clone(), pooled_prompt_embeds.clone(), seed)
        image_np = initial_image.detach().clone().cpu().numpy()
        image_np = (image_np * 255).astype(np.uint8)
        pil_image = Image.fromarray(image_np)
        pil_image.save(f"{results_folder}/it_0.png")

        initial_fitness, initial_aesthetic_score, initial_clip_score, initial_fitness_1, initial_fitness_2 = evaluate(trainable_params_init, seed, text_embeddings_init_shape, selected_prompt)

    time_list = [0]
    best_aesthetic_score_overall = initial_aesthetic_score
    best_clip_score_overall = initial_clip_score
    best_fitness_overall = initial_fitness
    best_text_embeddings_overall = trainable_params_init

    start_time = time.time()
    generation = 0

    max_fit_list = [-initial_fitness]
    avg_fit_list = [-initial_fitness]
    std_fit_list = [0]

    max_aesthetic_score_list = [initial_aesthetic_score]
    avg_aesthetic_score_list = [initial_aesthetic_score]
    std_aesthetic_score_list = [0]

    max_clip_score_list = [initial_clip_score]
    avg_clip_score_list = [initial_clip_score]
    std_clip_score_list = [0]

    max_fitness_1_list = [initial_fitness_1]
    avg_fitness_1_list = [initial_fitness_1]
    std_fitness_1_list = [0]

    max_fitness_2_list = [initial_fitness_2]
    avg_fitness_2_list = [initial_fitness_2]
    std_fitness_2_list = [0]

    while not es.stop():
        print(f"Generation {generation+1}/{NUM_GENERATIONS}")

        os.makedirs(results_folder+"/gen_%d" % (generation+1), exist_ok=True)

        # Ask for new candidate solutions
        solutions = es.ask()
        # Evaluate candidate solutions
        tmp_fitnesses = []
        tmp_fitness_1 = []
        tmp_fitness_2 = []
        aesthetic_scores = []
        clip_scores = []

        ind_id = 1
        for x in solutions:
            save_path = results_folder + "/gen_%d/id_%d.png" % (generation+1, ind_id)
            fitness, aesthetic_score, clip_score, fitness_1, fitness_2 = evaluate(x, seed, text_embeddings_init_shape, selected_prompt, save_path)
            tmp_fitnesses.append(fitness)
            tmp_fitness_1.append(fitness_1)
            tmp_fitness_2.append(fitness_2)
            aesthetic_scores.append(aesthetic_score)
            clip_scores.append(clip_score)
            ind_id += 1
        # Tell CMA-ES the fitnesses
        es.tell(solutions, tmp_fitnesses)

        # Record statistics
        fitnesses = [-f for f in tmp_fitnesses]  # Convert back to positive scores

        max_fit = max(fitnesses)
        avg_fit = np.mean(fitnesses)
        std_fit = np.std(fitnesses)

        max_fit_list.append(max_fit)
        avg_fit_list.append(avg_fit)
        std_fit_list.append(std_fit)

        max_aesthetic_score = max(aesthetic_scores)
        avg_aesthetic_score = np.mean(aesthetic_scores)
        std_aesthetic_score = np.std(aesthetic_scores)

        max_aesthetic_score_list.append(max_aesthetic_score)
        avg_aesthetic_score_list.append(avg_aesthetic_score)
        std_aesthetic_score_list.append(std_aesthetic_score)

        max_clip_score = max(clip_scores)
        avg_clip_score = np.mean(clip_scores)
        std_clip_score = np.std(clip_scores)

        max_clip_score_list.append(max_clip_score)
        avg_clip_score_list.append(avg_clip_score)
        std_clip_score_list.append(std_clip_score)

        max_fitness_1 = max(tmp_fitness_1)
        avg_fitness_1 = np.mean(tmp_fitness_1)
        std_fitness_1 = np.std(tmp_fitness_1)

        max_fitness_1_list.append(max_fitness_1)
        avg_fitness_1_list.append(avg_fitness_1)
        std_fitness_1_list.append(std_fitness_1)

        max_fitness_2 = max(tmp_fitness_2)
        avg_fitness_2 = np.mean(tmp_fitness_2)
        std_fitness_2 = np.std(tmp_fitness_2)

        max_fitness_2_list.append(max_fitness_2)
        avg_fitness_2_list.append(avg_fitness_2)
        std_fitness_2_list.append(std_fitness_2)

        # Get best solution so far
        best_x = es.result.xbest
        best_fitness = -es.result.fbest  # Convert back to positive score

        with torch.no_grad():
            # Generate and save the best image
            split = np.prod(text_embeddings_init_shape[0])
            best_pe  = torch.tensor(best_x[:split],  dtype=torch.float32, device=device).view(text_embeddings_init_shape[0])
            best_ppe = torch.tensor(best_x[split:], dtype=torch.float32, device=device).view(text_embeddings_init_shape[1])
            best_image = generate_image_from_embeddings(best_pe, best_ppe, seed)
            image_np = best_image.detach().clone().cpu().numpy()
            image_np = (image_np * 255).astype(np.uint8)
            pil_image = Image.fromarray(image_np)
            pil_image.save(results_folder + "/best_%d.png" % (generation+1))

        if best_fitness > best_fitness_overall:
            best_fitness_overall = best_fitness
            best_text_embeddings_overall = best_x

        generation += 1

        elapsed_time = time.time() - start_time
        generations_done = generation
        generations_left = NUM_GENERATIONS - generations_done
        average_time_per_generation = elapsed_time / generations_done
        estimated_time_remaining = average_time_per_generation * generations_left

        formatted_time_remaining = format_time(estimated_time_remaining)

        time_list.append(elapsed_time)

        # Save the metrics
        results = pd.DataFrame({
            "generation": list(range(0, generation + 1)),
            "prompt": [selected_prompt] + [''] * generation,
            "category": [category] + [''] * generation,
            "avg_fitness": avg_fit_list,
            "std_fitness": std_fit_list,
            "max_fitness": max_fit_list,
            #"avg_fitness_1": avg_fitness_1_list,
            #"std_fitness_1": std_fitness_1_list,
            #"max_fitness_1": max_fitness_1_list,
            #"avg_fitness_2": avg_fitness_2_list,
            #"std_fitness_2": std_fitness_2_list,
            #"max_fitness_2": max_fitness_2_list,
            "avg_aesthetic_score": avg_aesthetic_score_list,
            "std_aesthetic_score": std_aesthetic_score_list,
            "max_aesthetic_score": max_aesthetic_score_list,
            "avg_clip_score": avg_clip_score_list,
            "std_clip_score": std_clip_score_list,
            "max_clip_score": max_clip_score_list,
            "elapsed_time": time_list
        })

        results.to_csv(f"{results_folder}/fitness_results.csv", index=False, na_rep='nan')

        # # Convert embeddings to JSON strings
        # embeddings_json = [json.dumps(embedding.tolist()) for embedding in solutions]

        # results_generation = pd.DataFrame({
        #     "fitnesses": fitnesses,
        #     "embeddings": embeddings_json
        # })

        # results_generation.to_csv(f"{results_folder}/gen_{generation}/fitness_embeddings.csv", index=False, na_rep='nan')

        # Plot and save the fitness evolution
        save_plot_results(results, results_folder)

        # Print stats
        print(f"Seed {seed_number} Generation {generation}/{NUM_GENERATIONS}: Max fitness: {max_fit}, Avg fitness: {avg_fit}, Max aesthetic score: {max_aesthetic_score}, Avg aesthetic score: {avg_aesthetic_score}, Max clip score: {max_clip_score}, Avg clip score: {avg_clip_score}, Estimated time remaining: {formatted_time_remaining}")

    # Save the overall best image
    with torch.no_grad():
        split = np.prod(text_embeddings_init_shape[0])
        best_overall_pe  = torch.tensor(best_text_embeddings_overall[:split],  dtype=torch.float32, device=device).view(text_embeddings_init_shape[0])
        best_overall_ppe = torch.tensor(best_text_embeddings_overall[split:], dtype=torch.float32, device=device).view(text_embeddings_init_shape[1])
        best_image = generate_image_from_embeddings(best_overall_pe, best_overall_ppe, seed)
    best_image_np = best_image.detach().cpu().numpy()
    best_image_np = (best_image_np * 255).astype(np.uint8)
    pil_image = Image.fromarray(best_image_np)
    pil_image.save(f"{results_folder}/best_all.png")

def plot_mean_std(x_axis, m_vec, std_vec, description, title=None, y_label=None, x_label=None):
    lower_bound = [M_new - Sigma for M_new, Sigma in zip(m_vec, std_vec)]
    upper_bound = [M_new + Sigma for M_new, Sigma in zip(m_vec, std_vec)]

    plt.plot(x_axis, m_vec, '--', label=description + " Avg.")
    plt.fill_between(x_axis, lower_bound, upper_bound, alpha=.3, label=description + " Avg. ± SD")
    if title is not None:
        plt.title(title)
    if y_label is not None:
        plt.ylabel(y_label)
    if x_label is not None:
        plt.xlabel(x_label)

def save_plot_results(results, results_folder):
    # Plot main fitness evolution
    plt.figure(figsize=(10, 6))  # Increase figure size
    plot_mean_std(results['generation'], results['avg_fitness'], results['std_fitness'], "Fitness")
    plt.plot(results['generation'], results['max_fitness'], 'r-', label="Best Fitness")
    #plot_mean_std(results['generation'], results['avg_fitness_1'], results['std_fitness_1'], "F1 (Aesthetic Score)")
    #plt.plot(results['generation'], results['max_fitness_1'], 'orange', label="Best F1")
    #plot_mean_std(results['generation'], results['avg_fitness_2'], results['std_fitness_2'], "F2 (CLIP Score)")
    #plt.plot(results['generation'], results['max_fitness_2'], 'green', label="Best F2")
    plt.ylim(0, 1.1)
    plt.xlabel('Generation')
    plt.ylabel('Fitness')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))  # Move legend outside the plot
    plt.tight_layout()  # Adjust layout
    plt.savefig(results_folder + "/fitness_evolution.png")

    # Plot aesthetic score evolution
    plt.figure(figsize=(10, 6))  # Increase figure size
    plot_mean_std(results['generation'], results['avg_aesthetic_score'], results['std_aesthetic_score'], "Population")
    plt.plot(results['generation'], results['max_aesthetic_score'], 'r-', label="Best")
    plt.ylim(0, 10)
    plt.xlabel('Generation')
    plt.ylabel('Aesthetic Score')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))  # Move legend outside the plot
    plt.tight_layout()  # Adjust layout
    plt.savefig(results_folder + "/aesthetic_score_evolution.png")

    # Plot clip score evolution
    plt.figure(figsize=(10, 6))  # Increase figure size
    plot_mean_std(results['generation'], results['avg_clip_score'], results['std_clip_score'], "Population")
    plt.plot(results['generation'], results['max_clip_score'], 'r-', label="Best")
    plt.ylim(0, 0.6)
    plt.xlabel('Generation')
    plt.ylabel('CLIP Score')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))  # Move legend outside the plot
    plt.tight_layout()  # Adjust layout
    plt.savefig(results_folder + "/clip_score_evolution.png")

def aggregate_results():
    # Initialize aggregated_data as None
    aggregated_data = None

    # Iterate over all subdirectories
    for folder_name in os.listdir(OUTPUT_FOLDER):
        if folder_name.startswith(f"results_{model_name}_"):
            seed = folder_name.split("_")[-2]  # Extract the seed number
            prompt_number = folder_name.split("_")[-1]  # Extract the prompt number

            file_path = os.path.join(OUTPUT_FOLDER, folder_name, "fitness_results.csv")

            if os.path.exists(file_path):
                # Read the CSV file
                df = pd.read_csv(file_path)

                df = df.drop(columns=["prompt", "category"]) 

                df = df.rename(columns={"avg_fitness": f"avg_fitness_{seed}_{prompt_number}"})
                df = df.rename(columns={"max_fitness": f"max_fitness_{seed}_{prompt_number}"})
                df = df.rename(columns={"std_fitness": f"std_fitness_{seed}_{prompt_number}"})
                #df = df.rename(columns={"avg_fitness_1": f"avg_fitness_1_{seed}_{prompt_number}"})
                #df = df.rename(columns={"max_fitness_1": f"max_fitness_1_{seed}_{prompt_number}"})
                #df = df.rename(columns={"std_fitness_1": f"std_fitness_1_{seed}_{prompt_number}"})
                #df = df.rename(columns={"avg_fitness_2": f"avg_fitness_2_{seed}_{prompt_number}"})
                #df = df.rename(columns={"max_fitness_2": f"max_fitness_2_{seed}_{prompt_number}"})
                #df = df.rename(columns={"std_fitness_2": f"std_fitness_2_{seed}_{prompt_number}"})
                df = df.rename(columns={"avg_aesthetic_score": f"avg_aesthetic_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"max_aesthetic_score": f"max_aesthetic_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"std_aesthetic_score": f"std_aesthetic_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"avg_clip_score": f"avg_clip_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"max_clip_score": f"max_clip_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"std_clip_score": f"std_clip_score_{seed}_{prompt_number}"})
                df = df.rename(columns={"elapsed_time": f"elapsed_time_{seed}_{prompt_number}"})

                if aggregated_data is None:
                    aggregated_data = df

                aggregated_data = pd.merge(aggregated_data, df, on="generation", how="outer")
            else:
                print(f"File not found: {file_path}")

    # Ensure aggregated_data is not None before saving
    if aggregated_data is not None:
        # Save the aggregated data to an Excel file
        output_file = os.path.join(OUTPUT_FOLDER, "aggregated_score_results.xlsx")
        aggregated_data.to_excel(output_file, index=False)
        print(f"Aggregated results saved to {output_file}")
    else:
        print("No data was aggregated. Check the input folders and files.")

    data = pd.read_excel(output_file)

    # Calculate the average fitness across all seeds for each iteration
    data['avg_fitness'] = data.filter(like='avg_fitness_').mean(axis=1)
    # Calculate the standard deviation of fitness across all seeds for each iteration
    data['std_fitness'] = data.filter(like='std_fitness_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_avg_fitness'] = data.filter(like='max_fitness_').mean(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_std_fitness'] = data.filter(like='max_fitness_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['max_fitness'] = data.filter(like='max_fitness_').max(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['avg_aesthetic_score'] = data.filter(like='avg_aesthetic_score_').mean(axis=1)
    # Calculate the standard deviation of fitness across all seeds for each iteration
    data['std_aesthetic_score'] = data.filter(like='std_aesthetic_score_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_avg_aesthetic_score'] = data.filter(like='max_aesthetic_score_').mean(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_std_aesthetic_score'] = data.filter(like='max_aesthetic_score_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['max_aesthetic_score'] = data.filter(like='max_aesthetic_score_').max(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['avg_clip_score'] = data.filter(like='avg_clip_score_').mean(axis=1)
    # Calculate the standard deviation of fitness across all seeds for each iteration
    data['std_clip_score'] = data.filter(like='std_clip_score_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_avg_clip_score'] = data.filter(like='max_clip_score_').mean(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['best_std_clip_score'] = data.filter(like='max_clip_score_').std(axis=1)
    # Calculate the average fitness across all seeds for each iteration
    data['max_clip_score'] = data.filter(like='max_clip_score_').max(axis=1)

    # Fitness Evolution
    plt.figure(figsize=(10, 6))
    plot_mean_std(data['generation'], data['avg_fitness'], data['std_fitness'], "Population")
    plot_mean_std(data['generation'], data['best_avg_fitness'], data['best_std_fitness'], "Bests")
    plt.plot(data['generation'], data['max_fitness'], 'r-', label="Best")
    plt.ylim(0, 1.1)
    plt.xlabel('Generation')
    plt.ylabel('Fitness')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(OUTPUT_FOLDER + "/fitness_evolution.png")
    plt.close()

    # Aesthetic Score Evolution
    plt.figure(figsize=(10, 6))
    plot_mean_std(data['generation'], data['avg_aesthetic_score'], data['std_aesthetic_score'], "Population")
    plot_mean_std(data['generation'], data['best_avg_aesthetic_score'], data['best_std_aesthetic_score'], "Bests")
    plt.plot(data['generation'], data['max_aesthetic_score'], 'r-', label="Best")
    plt.ylim(0, 10.5)
    plt.xlabel('Generation')
    plt.ylabel('Aesthetic Score')
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(OUTPUT_FOLDER + "/aesthetic_score_evolution.png")
    plt.close()

    # CLIP Score Evolution
    plt.figure(figsize=(10, 6))
    plot_mean_std(data['generation'], data['avg_clip_score'], data['std_clip_score'], "Population")
    plot_mean_std(data['generation'], data['best_avg_clip_score'], data['best_std_clip_score'], "Bests")
    plt.plot(data['generation'], data['max_clip_score'], 'r-', label="Best")
    plt.xlabel('Generation')
    plt.ylabel('CLIP Score')
    plt.ylim(0, 0.6)
    plt.grid()
    plt.legend(loc="upper left", bbox_to_anchor=(1, 1))
    plt.tight_layout()
    plt.savefig(OUTPUT_FOLDER + "/clip_score_evolution.png")
    plt.close()

    # Initialize PowerPoint presentation
    presentation = Presentation()

    # Collect folders with seed numbers
    folders = []
    for folder_name in os.listdir(OUTPUT_FOLDER):
        if folder_name.startswith("results_"):
            # Extract the seed number from the folder name
            seed_number = int(folder_name.split("_")[-2])  # Convert to integer for sorting
            prompt_number = int(folder_name.split("_")[-1])  # Extract prompt number
            folder_path = os.path.join(OUTPUT_FOLDER, folder_name)
            folders.append((seed_number, prompt_number, folder_path))

    # Sort folders by prompt number in ascending order
    folders.sort(key=lambda x: x[1])

    # Iterate over sorted folders
    for seed_number, prompt_number, folder_path in folders:
        # Paths for required images and CSV file
        it_0_path = os.path.join(folder_path, "it_0.png")
        best_all_path = os.path.join(folder_path, "best_all.png")
        fitness_evolution_path = os.path.join(folder_path, "fitness_evolution.png")
        aesthetic_evolution_path = os.path.join(folder_path, "aesthetic_score_evolution.png")
        clip_evolution_path = os.path.join(folder_path, "clip_score_evolution.png")
        csv_path = os.path.join(folder_path, "fitness_results.csv")

        # Extract scores and prompt from CSV
        fitness_initial = None
        fitness_best = None
        aesthetic_initial = None
        aesthetic_best = None
        clip_initial = None
        clip_best = None
        prompt_text = None
        category = None

        if os.path.exists(csv_path):
            with open(csv_path, 'r') as csvfile:
                reader = csv.DictReader(csvfile)
                rows = list(reader)
                if rows:
                    # Initial values from the first row
                    first_row = rows[0]
                    fitness_initial = float(first_row['max_fitness'])
                    aesthetic_initial = float(first_row['max_aesthetic_score'])
                    clip_initial = float(first_row['max_clip_score'])
                    prompt_text = first_row['prompt']
                    category = first_row['category']

                    # Find the row with the best (maximum) max_fitness
                    best_row = max(rows, key=lambda r: float(r['max_fitness']))
                    fitness_best = float(best_row['max_fitness'])
                    aesthetic_best = float(best_row['max_aesthetic_score'])
                    clip_best = float(best_row['max_clip_score'])

        # Slide 1: it_0.png and it_1000.png
        if os.path.exists(it_0_path) and os.path.exists(best_all_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            # Add prompt below the title
            if prompt_text:
                left = Inches(0.5)
                top = Inches(1)
                width = Inches(9)
                textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
                textbox.text = f"Prompt: {prompt_text}"

            if category:
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
                textbox.text = f"Category: {category}"

            # Add it_0.png
            slide.shapes.add_picture(it_0_path, Inches(0.5), Inches(2), height=Inches(4))

            # Add legend below it_0.png
            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Initial iteration"
            if fitness_initial is not None:
                text += f"\nInitial Fitness: {fitness_initial:.4f}"
            if aesthetic_initial is not None:
                text += f"\nAesthetic Score: {aesthetic_initial:.4f}"
            if clip_initial is not None:
                text += f"\nCLIP Score: {clip_initial:.4f}"
            textbox.text = text

            # Add it_1000.png
            slide.shapes.add_picture(best_all_path, Inches(5.5), Inches(2), height=Inches(4))

            # Add legend below it_1000.png
            left = Inches(5.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Best iteration"
            if fitness_best is not None:
                text += f"\nBest fitness: {fitness_best:.4f}"
            if aesthetic_best is not None:
                text += f"\nAesthetic Score: {aesthetic_best:.4f}"
            if clip_best is not None:
                text += f"\nCLIP Score: {clip_best:.4f}"
            textbox.text = text

        # Slide 2: fitness_evolution
        if os.path.exists(fitness_evolution_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            # Add aesthetic_evolution.png
            slide.shapes.add_picture(fitness_evolution_path, Inches(0), Inches(2), height=Inches(4))

            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Fitness evolution"
            textbox.text = text

        # Slide 3: clip_score_evolution.png and aesthetic_score_evolution.png
        if os.path.exists(clip_evolution_path) and os.path.exists(aesthetic_evolution_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            # Add aesthetic_evolution.png
            slide.shapes.add_picture(clip_evolution_path, Inches(0), Inches(2), height=Inches(4))

            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "CLIP score evolution"
            textbox.text = text

            # Add loss_evolution.png
            slide.shapes.add_picture(aesthetic_evolution_path, Inches(5), Inches(2), height=Inches(4))

            left = Inches(5.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Aesthetic score evolution"
            textbox.text = text

    output_filename = os.path.join(OUTPUT_FOLDER, f"summary.pptx")
    # Save the presentation
    presentation.save(output_filename)
    print(f"Presentation saved as {output_filename}")

if __name__ == "__main__":
    seed_number = 1
    for seed in seed_list:
        prompt_number = 1
        for prompt, category in selected_prompts_with_category:
            print(f"Running seed {seed}, prompt: {prompt} (Category: {category})")
            main(seed, seed_number, prompt, category, prompt_number)
            print(f"Run with seed {seed} and prompt '{prompt}' finished!")
            aggregate_results()
            prompt_number += 1
        seed_number += 1
    